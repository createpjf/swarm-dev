"""
core/tools.py
Built-in tool registry — OpenClaw-inspired agent tool system.

Architecture:
  - Tools are callable functions with JSON-schema parameters
  - Agent system prompts include tool descriptions
  - Agents invoke tools via structured JSON blocks in their output
  - Tool results are fed back to the agent as context

Tool categories (37 tools across 10 groups):
  - Web:        web_search (Brave + Perplexity), web_fetch (text + markdown)
  - Filesystem: read_file, write_file, edit_file, list_dir
  - Memory:     memory_search, memory_save, kb_search, kb_write
  - Task:       task_create, task_status, spawn_subagent
  - Automation: exec, cron, process
  - Skill:      check_skill_deps, install_skill_cli, search_skills, install_remote_skill
  - Browser:    browser_navigate, browser_click, browser_fill, browser_get_text,
                browser_screenshot, browser_evaluate, browser_page_info
  - Media:      screenshot, notify, analyze_image
  - Messaging:  send_mail, send_file, message
  - A2A:        a2a_delegate (delegate to external agents via A2A protocol)

Access control:
  - Tool profiles: "minimal", "coding", "full"
  - Per-agent allow/deny lists in agents.yaml
  - Deny always wins over allow

Usage in agents.yaml:
  agents:
    - id: researcher
      tools:
        profile: "full"            # or "coding", "minimal"
        allow: ["web_search"]      # additional allowlist
        deny: ["exec", "write_file"]  # explicit denylist
"""

from __future__ import annotations

import json
import logging
import os
import re
import subprocess
import tempfile
import time
import urllib.parse
import urllib.request
from datetime import datetime, timezone
from typing import Any, Callable, Optional

logger = logging.getLogger(__name__)

# ── Audit logging for sensitive tool calls ──
_AUDIT_LOG = ".logs/tool_audit.log"


def _audit_log(tool_name: str, agent_id: str = "unknown", **details):
    """Append an audit entry for sensitive tool invocations."""
    try:
        os.makedirs(os.path.dirname(_AUDIT_LOG), exist_ok=True)
        entry = {
            "ts": datetime.now(timezone.utc).isoformat(),
            "tool": tool_name,
            "agent": agent_id,
        }
        entry.update(details)
        with open(_AUDIT_LOG, "a") as f:
            f.write(json.dumps(entry) + "\n")
    except OSError:
        pass


def _is_allowed_path(abs_path: str) -> bool:
    """Check if path is within project directory or a temp directory."""
    cwd = os.path.abspath(".")
    real_path = os.path.realpath(abs_path)
    # System temp dir (macOS: /private/var/folders/.../T, Linux: /tmp)
    sys_tmp = os.path.realpath(tempfile.gettempdir())
    # Also allow /tmp/ explicitly (macOS symlinks /tmp → /private/tmp)
    slash_tmp = os.path.realpath("/tmp")
    return (real_path.startswith(cwd)
            or real_path.startswith(sys_tmp + os.sep)
            or real_path.startswith(slash_tmp + os.sep))


# ══════════════════════════════════════════════════════════════════════════════
#  TOOL SCHEMA
# ══════════════════════════════════════════════════════════════════════════════

class Tool:
    """A single tool that agents can invoke."""

    def __init__(self, name: str, description: str,
                 parameters: dict[str, dict],
                 handler: Callable[..., dict],
                 group: str = "",
                 requires_env: list[str] | None = None):
        self.name = name
        self.description = description
        self.parameters = parameters        # {param_name: {type, description, required?}}
        self.handler = handler
        self.group = group                  # e.g. "web", "automation", "media", "fs"
        self.requires_env = requires_env or []

    def is_available(self) -> bool:
        """Check if required env vars are set."""
        for env in self.requires_env:
            if not os.environ.get(env):
                return False
        return True

    def to_prompt(self) -> str:
        """Generate tool description for system prompt injection."""
        params_desc = []
        example_params = {}
        for pname, pinfo in self.parameters.items():
            req = " (required)" if pinfo.get("required") else ""
            params_desc.append(
                f"    - {pname}: {pinfo.get('type', 'string')} — "
                f"{pinfo.get('description', '')}{req}")
            if pinfo.get("required"):
                example_params[pname] = f"<{pname}>"
        params_str = "\n".join(params_desc) if params_desc else "    (no parameters)"
        example_json = json.dumps(
            {"tool": self.name, "params": example_params}, ensure_ascii=False)
        return (f"### {self.name}\n"
                f"{self.description}\n"
                f"  Parameters:\n{params_str}\n"
                f"  Example: {example_json}\n")

    def to_schema(self) -> dict:
        """Generate JSON schema for function-calling LLMs."""
        properties = {}
        required = []
        for pname, pinfo in self.parameters.items():
            properties[pname] = {
                "type": pinfo.get("type", "string"),
                "description": pinfo.get("description", ""),
            }
            if pinfo.get("required"):
                required.append(pname)
        return {
            "name": self.name,
            "description": self.description,
            "parameters": {
                "type": "object",
                "properties": properties,
                "required": required,
            },
        }

    def execute(self, **kwargs) -> dict:
        """Execute the tool handler."""
        try:
            return self.handler(**kwargs)
        except Exception as e:
            logger.error("Tool %s error: %s", self.name, e)
            return {"ok": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════════
#  TOOL PROFILES (access control)
# ══════════════════════════════════════════════════════════════════════════════

TOOL_PROFILES = {
    "minimal": {"web_search", "web_fetch", "memory_search", "kb_search",
                "check_skill_deps", "install_skill_cli",
                "search_skills", "install_remote_skill"},
    "coding": {"web_search", "web_fetch", "exec", "read_file", "write_file",
               "edit_file", "list_dir", "generate_doc",
               "process", "cron_list", "cron_add",
               "notify", "transcribe", "tts", "list_voices",
               "memory_search", "memory_save",
               "kb_search", "kb_write", "task_create", "task_status",
               "send_mail", "send_file", "message", "check_skill_deps", "install_skill_cli",
               "search_skills", "install_remote_skill",
               "browser_navigate", "browser_click", "browser_fill",
               "browser_get_text", "browser_screenshot",
               "browser_evaluate", "browser_page_info",
               "a2a_delegate"},  # Phase 5: A2A delegation
    "full": None,  # None = all tools allowed
}

# Tool groups for bulk allow/deny
TOOL_GROUPS = {
    "group:web": ["web_search", "web_fetch"],
    "group:automation": ["exec", "cron_list", "cron_add", "process"],
    "group:media": ["screenshot", "notify", "transcribe", "tts", "list_voices",
                    "analyze_image"],
    "group:fs": ["read_file", "write_file", "edit_file", "list_dir",
                 "generate_doc", "workspace_status"],
    "group:memory": ["memory_search", "memory_save", "kb_search", "kb_write"],
    "group:task": ["task_create", "task_status", "spawn_subagent"],
    "group:skill": ["check_skill_deps", "install_skill_cli",
                    "search_skills", "install_remote_skill"],
    "group:messaging": ["send_mail", "send_file", "message"],
    "group:browser": ["browser_navigate", "browser_click", "browser_fill",
                      "browser_get_text", "browser_screenshot",
                      "browser_evaluate", "browser_page_info"],
    "group:a2a": ["a2a_delegate"],
}


# ══════════════════════════════════════════════════════════════════════════════
#  BUILT-IN TOOL HANDLERS
# ══════════════════════════════════════════════════════════════════════════════

# ── Web tool cache (15-minute TTL) ──
_web_cache: dict[str, tuple[float, dict]] = {}
_WEB_CACHE_TTL = 900  # 15 minutes


def _cache_get(key: str) -> dict | None:
    """Get cached result if still fresh."""
    if key in _web_cache:
        ts, result = _web_cache[key]
        if time.time() - ts < _WEB_CACHE_TTL:
            result["_cached"] = True
            return result
        del _web_cache[key]
    return None


def _cache_set(key: str, result: dict) -> dict:
    """Cache a result and evict stale entries (max 100)."""
    now = time.time()
    # Evict stale
    stale = [k for k, (ts, _) in _web_cache.items() if now - ts >= _WEB_CACHE_TTL]
    for k in stale:
        del _web_cache[k]
    # Evict oldest if over limit
    if len(_web_cache) >= 100:
        oldest_key = min(_web_cache, key=lambda k: _web_cache[k][0])
        del _web_cache[oldest_key]
    _web_cache[key] = (now, result)
    return result


def _is_private_hostname(hostname: str) -> bool:
    """Block requests to private/internal hostnames."""
    import socket
    blocked = {"localhost", "127.0.0.1", "0.0.0.0", "::1",
               "metadata.google.internal", "169.254.169.254"}
    if hostname.lower() in blocked:
        return True
    # Block private IP ranges
    try:
        ip = socket.gethostbyname(hostname)
        parts = ip.split(".")
        if len(parts) == 4:
            a, b = int(parts[0]), int(parts[1])
            if a == 10 or (a == 172 and 16 <= b <= 31) or (a == 192 and b == 168):
                return True
            if ip.startswith("127.") or ip.startswith("0."):
                return True
    except (socket.gaierror, ValueError):
        pass
    return False


def _handle_web_search(query: str, count: int = 5, freshness: str = "",
                       country: str = "", search_lang: str = "",
                       ui_lang: str = "", provider: str = "", **_) -> dict:
    """Search the web. Supports Brave, Perplexity Sonar, and Kimi/Moonshot.

    Provider auto-detection:
      - If BRAVE_API_KEY is set → Brave Search (default)
      - If PERPLEXITY_API_KEY is set → Perplexity Sonar (fallback or explicit)
      - If MOONSHOT_API_KEY is set → Kimi search (best for Chinese queries)
      - Use provider="perplexity"|"kimi" to force a specific provider
    """
    # Check cache
    cache_key = f"search:{query}:{count}:{freshness}:{country}:{search_lang}:{provider}"
    cached = _cache_get(cache_key)
    if cached:
        return cached

    brave_key = os.environ.get("BRAVE_API_KEY", "")
    pplx_key = os.environ.get("PERPLEXITY_API_KEY", "")
    kimi_key = os.environ.get("MOONSHOT_API_KEY", "")
    use_perplexity = (provider == "perplexity" and pplx_key) or (not brave_key and pplx_key and not kimi_key)
    use_kimi = (provider == "kimi" and kimi_key)

    if use_kimi:
        return _cache_set(cache_key,
                          _search_kimi(query, int(count), kimi_key))

    if use_perplexity:
        return _cache_set(cache_key,
                          _search_perplexity(query, int(count), pplx_key))

    if not brave_key:
        return {"ok": False, "error": "No search API configured. Set BRAVE_API_KEY "
                "(https://brave.com/search/api/) or PERPLEXITY_API_KEY."}

    params: dict[str, Any] = {"q": query, "count": min(int(count), 20)}
    if freshness:
        params["freshness"] = freshness
    if country:
        params["country"] = country        # e.g. "US", "CN", "JP"
    if search_lang:
        params["search_lang"] = search_lang  # e.g. "en", "zh", "ja"
    if ui_lang:
        params["ui_lang"] = ui_lang          # e.g. "en-US", "zh-CN"

    url = ("https://api.search.brave.com/res/v1/web/search?"
           + urllib.parse.urlencode(params))
    req = urllib.request.Request(url, headers={
        "Accept": "application/json",
        "Accept-Encoding": "gzip",
        "X-Subscription-Token": brave_key,
    })

    try:
        with urllib.request.urlopen(req, timeout=10) as resp:
            raw = resp.read()
            if resp.headers.get("Content-Encoding") == "gzip":
                import gzip
                raw = gzip.decompress(raw)
            data = json.loads(raw)

        results = []
        for item in (data.get("web", {}).get("results", []))[:int(count)]:
            results.append({
                "title": item.get("title", ""),
                "url": item.get("url", ""),
                "snippet": item.get("description", ""),
            })
        result = {"ok": True, "query": query, "results": results,
                  "total": len(results), "provider": "brave"}
        return _cache_set(cache_key, result)
    except Exception as e:
        # Auto-fallback: try Kimi, then Perplexity if Brave fails
        if kimi_key and provider != "brave":
            logger.warning("Brave search failed, falling back to Kimi: %s", e)
            return _cache_set(cache_key,
                              _search_kimi(query, int(count), kimi_key))
        if pplx_key and provider != "brave":
            logger.warning("Brave search failed, falling back to Perplexity: %s", e)
            return _cache_set(cache_key,
                              _search_perplexity(query, int(count), pplx_key))
        return {"ok": False, "error": f"Search failed: {e}"}


def _search_perplexity(query: str, count: int, api_key: str) -> dict:
    """Search using Perplexity Sonar API (chat-based search)."""
    payload = json.dumps({
        "model": "sonar",
        "messages": [{"role": "user", "content": query}],
        "max_tokens": 1024,
    }).encode()

    req = urllib.request.Request(
        "https://api.perplexity.ai/chat/completions",
        data=payload,
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
            "Accept": "application/json",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(req, timeout=15) as resp:
            data = json.loads(resp.read())

        content = data.get("choices", [{}])[0].get("message", {}).get("content", "")
        citations = data.get("citations", [])

        results = []
        for i, cite in enumerate(citations[:count]):
            results.append({
                "title": cite if isinstance(cite, str) else cite.get("title", f"Source {i+1}"),
                "url": cite if isinstance(cite, str) else cite.get("url", ""),
                "snippet": "",
            })
        # If no structured citations, return the answer text as a single result
        if not results:
            results.append({
                "title": "Perplexity Answer",
                "url": "",
                "snippet": content[:500],
            })

        return {"ok": True, "query": query, "results": results,
                "total": len(results), "provider": "perplexity",
                "answer": content[:2000]}
    except Exception as e:
        return {"ok": False, "error": f"Perplexity search failed: {e}"}


def _search_kimi(query: str, count: int, api_key: str) -> dict:
    """Search using Moonshot/Kimi API with built-in $web_search tool.

    Optimized for Chinese-language queries. Uses moonshot-v1-auto model
    with the $web_search builtin tool to get grounded search results.
    """
    payload = json.dumps({
        "model": "moonshot-v1-auto",
        "messages": [
            {"role": "system",
             "content": ("You are a search assistant. Use web_search to find "
                         "relevant information, then provide a concise summary "
                         "with source URLs. Reply in the same language as the query.")},
            {"role": "user", "content": query},
        ],
        "tools": [{
            "type": "builtin_function",
            "function": {"name": "$web_search"},
        }],
        "max_tokens": 2048,
    }).encode()

    req = urllib.request.Request(
        "https://api.moonshot.cn/v1/chat/completions",
        data=payload,
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
            "Accept": "application/json",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(req, timeout=20) as resp:
            data = json.loads(resp.read())

        content = (data.get("choices", [{}])[0]
                   .get("message", {}).get("content", ""))

        # Extract URLs from the response content
        import re as _re
        urls = _re.findall(r'https?://[^\s\)\"\'<>]+', content)
        results = []
        for i, url in enumerate(urls[:count]):
            results.append({
                "title": f"Source {i+1}",
                "url": url,
                "snippet": "",
            })
        if not results:
            results.append({
                "title": "Kimi Answer",
                "url": "",
                "snippet": content[:500],
            })

        return {"ok": True, "query": query, "results": results,
                "total": len(results), "provider": "kimi",
                "answer": content[:2000]}
    except Exception as e:
        return {"ok": False, "error": f"Kimi search failed: {e}"}


def _handle_web_fetch(url: str, max_chars: int = 8000,
                      extract_mode: str = "text",
                      timeout: int = 15, **_) -> dict:
    """Fetch URL content and extract readable content.

    extract_mode:
      - "text" (default): Plain text extraction — removes all HTML tags
      - "markdown": Convert HTML to simplified Markdown (headings, links, lists)
    """
    # Validate URL
    parsed = urllib.parse.urlparse(url)
    if not parsed.scheme or not parsed.netloc:
        return {"ok": False, "error": "Invalid URL — must include scheme (https://)"}

    # Block private hostnames
    hostname = parsed.hostname or ""
    if _is_private_hostname(hostname):
        return {"ok": False, "error": f"Blocked: private/internal hostname '{hostname}'"}

    # Check cache
    cache_key = f"fetch:{url}:{extract_mode}:{max_chars}"
    cached = _cache_get(cache_key)
    if cached:
        return cached

    try:
        req = urllib.request.Request(url, headers={
            "User-Agent": "CleoBot/1.0 (https://github.com/createpjf/cleo-dev)",
            "Accept-Encoding": "gzip, deflate",
            "Accept": "text/html,application/xhtml+xml,text/plain,*/*",
        })
        # Follow up to 5 redirects (urllib default), but cap response size
        with urllib.request.urlopen(req, timeout=int(timeout)) as resp:
            content_type = resp.headers.get("Content-Type", "")
            encoding = resp.headers.get("Content-Encoding", "")
            # Cap at 1MB to prevent memory issues
            raw = resp.read(1_000_000)
            final_url = resp.url  # capture redirect target

        # Decompress if gzipped
        if encoding == "gzip":
            import gzip
            raw = gzip.decompress(raw)
        elif encoding == "deflate":
            import zlib
            raw = zlib.decompress(raw)

        # Try to detect charset from content-type
        charset = "utf-8"
        if "charset=" in content_type.lower():
            charset = content_type.lower().split("charset=")[-1].split(";")[0].strip()
        text = raw.decode(charset, errors="ignore")

        if "html" in content_type.lower():
            if extract_mode == "markdown":
                text = _html_to_markdown(text)
            else:
                text = _html_to_text(text)

        text = text[:int(max_chars)]
        result = {"ok": True, "url": final_url or url, "content": text,
                  "chars": len(text), "extract_mode": extract_mode}
        if final_url and final_url != url:
            result["redirected_from"] = url
        return _cache_set(cache_key, result)
    except Exception as e:
        return {"ok": False, "error": f"Fetch failed: {e}"}


def _html_to_text(html: str) -> str:
    """Convert HTML to plain text — strips all tags."""
    # Remove script/style/nav/footer
    for tag in ("script", "style", "nav", "footer", "header", "noscript"):
        html = re.sub(rf'<{tag}[^>]*>.*?</{tag}>', '', html,
                       flags=re.DOTALL | re.IGNORECASE)
    # Remove comments
    html = re.sub(r'<!--.*?-->', '', html, flags=re.DOTALL)
    # Remove tags
    html = re.sub(r'<[^>]+>', ' ', html)
    # Decode HTML entities
    html = html.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
    html = html.replace("&quot;", '"').replace("&#39;", "'").replace("&nbsp;", " ")
    # Collapse whitespace
    html = re.sub(r'\s+', ' ', html).strip()
    return html


def _html_to_markdown(html: str) -> str:
    """Convert HTML to simplified Markdown — preserves structure."""
    # Remove script/style/nav/footer/noscript
    for tag in ("script", "style", "nav", "footer", "noscript"):
        html = re.sub(rf'<{tag}[^>]*>.*?</{tag}>', '', html,
                       flags=re.DOTALL | re.IGNORECASE)
    # Remove comments
    html = re.sub(r'<!--.*?-->', '', html, flags=re.DOTALL)

    # ── 1. Inline elements first (before block elements) ──

    # Convert bold/strong (use [\s>] boundary to avoid matching <body> etc.)
    html = re.sub(r'<(?:b|strong)(?:\s[^>]*)?>(.+?)</(?:b|strong)>', r'**\1**', html,
                   flags=re.DOTALL | re.IGNORECASE)
    # Convert italic/em (use [\s>] boundary to avoid matching <iframe> etc.)
    html = re.sub(r'<(?:i|em)(?:\s[^>]*)?>(.+?)</(?:i|em)>', r'*\1*', html,
                   flags=re.DOTALL | re.IGNORECASE)
    # Convert inline code
    html = re.sub(r'<code[^>]*>(.*?)</code>', r'`\1`', html,
                   flags=re.DOTALL | re.IGNORECASE)
    # Convert links: <a href="url">text</a> → [text](url)
    html = re.sub(
        r'<a[^>]*href="([^"]*)"[^>]*>(.*?)</a>',
        r'[\2](\1)', html, flags=re.DOTALL | re.IGNORECASE)
    # Convert images: <img src="url" alt="text"> → ![text](url)
    html = re.sub(r'<img[^>]*src="([^"]*)"[^>]*alt="([^"]*)"[^>]*/?>',
                   r'![\2](\1)', html, flags=re.IGNORECASE)
    html = re.sub(r'<img[^>]*src="([^"]*)"[^>]*/?>',
                   r'![image](\1)', html, flags=re.IGNORECASE)

    # ── 2. Block elements ──

    # Convert pre/code blocks (before heading/paragraph stripping)
    html = re.sub(r'<pre[^>]*>(.*?)</pre>', r'\n```\n\1\n```\n', html,
                   flags=re.DOTALL | re.IGNORECASE)

    # Convert headings
    for level in range(1, 7):
        prefix = "#" * level
        html = re.sub(
            rf'<h{level}[^>]*>(.*?)</h{level}>',
            rf'\n\n{prefix} \1\n\n', html,
            flags=re.DOTALL | re.IGNORECASE)

    # Convert list items
    html = re.sub(r'<li[^>]*>(.*?)</li>', r'\n- \1', html,
                   flags=re.DOTALL | re.IGNORECASE)

    # Convert paragraphs and line breaks
    html = re.sub(r'<br\s*/?>', '\n', html, flags=re.IGNORECASE)
    html = re.sub(r'<p[^>]*>(.*?)</p>', r'\n\n\1\n\n', html,
                   flags=re.DOTALL | re.IGNORECASE)

    # ── 3. Clean up ──

    # Remove remaining tags
    html = re.sub(r'<[^>]+>', '', html)
    # Decode HTML entities
    html = html.replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
    html = html.replace("&quot;", '"').replace("&#39;", "'").replace("&nbsp;", " ")
    # Collapse excessive newlines
    html = re.sub(r'\n{3,}', '\n\n', html)
    return html.strip()


def _handle_exec(command: str, timeout: int = 120, **_) -> dict:
    """Execute a shell command with approval gating."""
    from core.exec_tool import execute
    return execute(command=command, agent_id="tool", timeout=int(timeout))


def _handle_cron_list(**_) -> dict:
    """List all scheduled cron jobs."""
    from core.cron import list_jobs
    jobs = list_jobs()
    return {"ok": True, "jobs": jobs, "total": len(jobs)}


def _handle_cron_add(name: str, action: str, payload: str,
                     schedule_type: str, schedule: str, **_) -> dict:
    """Create a new scheduled job."""
    from core.cron import add_job
    job = add_job(name, action, payload, schedule_type, schedule)
    return {"ok": True, "job": job}


def _handle_process_list(**_) -> dict:
    """List running background processes."""
    try:
        result = subprocess.run(
            ["ps", "aux"], capture_output=True, text=True, timeout=10)
        lines = result.stdout.strip().split("\n")
        return {"ok": True, "processes": lines[:50],
                "total": len(lines) - 1}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_screenshot(**_) -> dict:
    """Take a screenshot of the current desktop (macOS)."""
    import platform
    if platform.system() != "Darwin":
        return {"ok": False, "error": "Screenshots only supported on macOS"}

    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = os.path.join("memory", f"screenshot_{ts}.png")
    os.makedirs("memory", exist_ok=True)

    try:
        result = subprocess.run(
            ["screencapture", "-x", path],
            capture_output=True, text=True, timeout=10)
        if result.returncode == 0 and os.path.exists(path):
            size = os.path.getsize(path)
            return {"ok": True, "path": path, "size": size}
        return {"ok": False, "error": result.stderr or "Failed to capture"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_notify(title: str, message: str = "", **_) -> dict:
    """Send a macOS notification."""
    import platform
    if platform.system() != "Darwin":
        return {"ok": False, "error": "Notifications only supported on macOS"}

    script = (f'display notification "{message}" '
              f'with title "{title}" sound name "Glass"')
    try:
        result = subprocess.run(
            ["osascript", "-e", script],
            capture_output=True, text=True, timeout=10)
        return {"ok": result.returncode == 0,
                "error": result.stderr if result.returncode else None}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_transcribe(file_path: str, language: str = "", **_) -> dict:
    """Transcribe audio using OpenAI Whisper API.

    Supports: mp3, mp4, mpeg, mpga, m4a, wav, webm, ogg, oga, flac.
    Requires OPENAI_API_KEY environment variable.
    """
    import httpx

    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        return {"ok": False, "error": "OPENAI_API_KEY not set — required for Whisper transcription"}

    abs_path = os.path.abspath(file_path)
    if not os.path.isfile(abs_path):
        return {"ok": False, "error": f"File not found: {file_path}"}

    # Validate extension
    ext = os.path.splitext(abs_path)[1].lower()
    allowed = {".mp3", ".mp4", ".mpeg", ".mpga", ".m4a", ".wav",
               ".webm", ".ogg", ".oga", ".flac"}
    if ext not in allowed:
        return {"ok": False,
                "error": f"Unsupported audio format '{ext}'. Supported: {', '.join(sorted(allowed))}"}

    # Check file size (Whisper limit: 25 MB)
    size = os.path.getsize(abs_path)
    if size > 25 * 1024 * 1024:
        return {"ok": False, "error": f"File too large ({size / 1024 / 1024:.1f} MB). Whisper limit is 25 MB."}

    base_url = os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1")
    url = f"{base_url}/audio/transcriptions"

    try:
        with open(abs_path, "rb") as f:
            files = {"file": (os.path.basename(abs_path), f)}
            data = {"model": "whisper-1"}
            if language:
                data["language"] = language

            resp = httpx.post(
                url,
                headers={"Authorization": f"Bearer {api_key}"},
                files=files,
                data=data,
                timeout=120,
            )

        if resp.status_code != 200:
            return {"ok": False, "error": f"Whisper API error {resp.status_code}: {resp.text[:300]}"}

        result = resp.json()
        text = result.get("text", "").strip()
        return {"ok": True, "text": text, "file": os.path.basename(abs_path),
                "size_kb": round(size / 1024, 1)}

    except httpx.TimeoutException:
        return {"ok": False, "error": "Whisper API request timed out (120s)"}
    except Exception as e:
        return {"ok": False, "error": f"Transcription failed: {e}"}


def _handle_tts(text: str, voice: str = "", speed: float = 1.0,
                provider: str = "", output_format: str = "mp3", **_) -> dict:
    """Synthesize text to speech audio file.

    Multi-provider TTS with automatic fallback:
      - OpenAI TTS (alloy/echo/nova/shimmer voices, needs OPENAI_API_KEY)
      - ElevenLabs (rachel/adam/sam voices, needs ELEVENLABS_API_KEY)
      - MiniMax (Chinese-optimized, needs MINIMAX_API_KEY + GROUP_ID)
      - Local (piper/sherpa-onnx, zero cost, needs binary installed)

    Returns path to generated audio file.
    """
    import asyncio
    try:
        from adapters.voice.tts_engine import get_tts_engine
    except ImportError:
        return {"ok": False, "error": "TTS engine not available"}

    engine = get_tts_engine()

    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor() as pool:
                result = pool.submit(
                    lambda: asyncio.run(engine.synthesize(
                        text, voice=voice, speed=speed,
                        output_format=output_format, provider=provider))
                ).result(timeout=90)
        else:
            result = loop.run_until_complete(engine.synthesize(
                text, voice=voice, speed=speed,
                output_format=output_format, provider=provider))
    except RuntimeError:
        result = asyncio.run(engine.synthesize(
            text, voice=voice, speed=speed,
            output_format=output_format, provider=provider))

    return result


def _handle_list_voices(provider: str = "", **_) -> dict:
    """List available TTS voices across all providers."""
    try:
        from adapters.voice.tts_engine import get_tts_engine
    except ImportError:
        return {"ok": False, "error": "TTS engine not available"}

    engine = get_tts_engine()
    return {
        "ok": True,
        "providers": engine.list_providers(),
        "voices": engine.list_voices(provider=provider),
    }


def _handle_read_file(path: str, max_lines: int = 200, **_) -> dict:
    """Read a file from the project directory."""
    abs_path = os.path.abspath(path)
    if not _is_allowed_path(abs_path):
        return {"ok": False, "error": "Cannot read files outside project or temp directory"}

    if not os.path.exists(abs_path):
        return {"ok": False, "error": f"File not found: {path}"}

    try:
        with open(abs_path, "r", encoding="utf-8", errors="ignore") as f:
            lines = f.readlines()[:int(max_lines)]
        content = "".join(lines)
        return {"ok": True, "path": path, "content": content,
                "lines": len(lines)}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_write_file(path: str, content: str, **kwargs) -> dict:
    """Write content to a file in the project directory."""
    abs_path = os.path.abspath(path)
    if not _is_allowed_path(abs_path):
        return {"ok": False, "error": "Cannot write files outside project or temp directory"}

    agent_id = kwargs.get("_agent_id", "unknown")
    _audit_log("write_file", agent_id=agent_id, path=path, size=len(content))

    try:
        os.makedirs(os.path.dirname(abs_path), exist_ok=True)
        with open(abs_path, "w", encoding="utf-8") as f:
            f.write(content)

        # Save agent metadata for collaboration tracking
        try:
            meta_path = abs_path + ".meta"
            with open(meta_path, "w") as mf:
                json.dump({
                    "agent": agent_id,
                    "task_id": kwargs.get("_task_id", ""),
                    "ts": time.time(),
                }, mf)
        except OSError:
            pass  # Metadata is optional

        return {"ok": True, "path": path, "size": len(content)}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_list_dir(path: str = ".", **_) -> dict:
    """List directory contents."""
    abs_path = os.path.abspath(path)
    if not _is_allowed_path(abs_path):
        return {"ok": False, "error": "Cannot list outside project or temp directory"}

    if not os.path.isdir(abs_path):
        return {"ok": False, "error": f"Not a directory: {path}"}

    try:
        entries = []
        for name in sorted(os.listdir(abs_path)):
            full = os.path.join(abs_path, name)
            entries.append({
                "name": name,
                "type": "dir" if os.path.isdir(full) else "file",
                "size": os.path.getsize(full) if os.path.isfile(full) else 0,
            })
        return {"ok": True, "path": path, "entries": entries[:100],
                "total": len(entries)}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_edit_file(path: str, old_str: str, new_str: str, **kwargs) -> dict:
    """Find-and-replace edit in a file (safe, project-scoped)."""
    abs_path = os.path.abspath(path)
    if not _is_allowed_path(abs_path):
        return {"ok": False, "error": "Cannot edit files outside project or temp directory"}

    _audit_log("edit_file", agent_id=kwargs.get("_agent_id", "unknown"),
               path=path)

    if not os.path.exists(abs_path):
        return {"ok": False, "error": f"File not found: {path}"}

    try:
        with open(abs_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()

        count = content.count(old_str)
        if count == 0:
            return {"ok": False, "error": "old_str not found in file"}
        if count > 1:
            return {"ok": False, "error": f"old_str found {count} times — must be unique (include more context)"}

        new_content = content.replace(old_str, new_str, 1)
        with open(abs_path, "w", encoding="utf-8") as f:
            f.write(new_content)

        return {"ok": True, "path": path, "replacements": 1}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_memory_search(query: str, limit: int = 5, agent_id: str = "tool", **_) -> dict:
    """Search episodic memory for past cases (problem→solution pairs)."""
    try:
        from adapters.memory.episodic import EpisodicMemory
        mem = EpisodicMemory(agent_id=agent_id)
        cases = mem.search_cases(query, limit=int(limit))
        return {"ok": True, "results": cases, "total": len(cases)}
    except ImportError:
        return {"ok": False, "error": "Episodic memory module not available"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_memory_save(problem: str, solution: str,
                        tags: str = "", agent_id: str = "tool", **_) -> dict:
    """Save a problem→solution case to episodic memory."""
    try:
        from adapters.memory.episodic import EpisodicMemory
        mem = EpisodicMemory(agent_id=agent_id)
        tag_list = [t.strip() for t in tags.split(",") if t.strip()] if tags else None
        case_id = mem.save_case(problem, solution, tags=tag_list)
        return {"ok": True, "case_id": case_id}
    except ImportError:
        return {"ok": False, "error": "Episodic memory module not available"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_kb_search(query: str, limit: int = 5, **_) -> dict:
    """Search the shared knowledge base for notes."""
    try:
        from adapters.memory.knowledge_base import KnowledgeBase
        kb = KnowledgeBase()
        notes = kb.search_notes(query, limit=int(limit))
        # Trim content for return
        results = []
        for n in notes:
            results.append({
                "topic": n.get("topic", ""),
                "slug": n.get("slug", ""),
                "content": n.get("content", "")[:500],
                "tags": n.get("tags", []),
                "author": n.get("author", ""),
            })
        return {"ok": True, "results": results, "total": len(results)}
    except ImportError:
        return {"ok": False, "error": "Knowledge base module not available"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_kb_write(topic: str, content: str, tags: str = "",
                     agent_id: str = "tool", **_) -> dict:
    """Create or update a note in the shared knowledge base."""
    try:
        from adapters.memory.knowledge_base import KnowledgeBase
        kb = KnowledgeBase()
        tag_list = [t.strip() for t in tags.split(",") if t.strip()] if tags else None
        slug = kb.create_note(topic, content, tags=tag_list, author=agent_id)
        return {"ok": True, "slug": slug, "topic": topic}
    except ImportError:
        return {"ok": False, "error": "Knowledge base module not available"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_task_create(description: str, **_) -> dict:
    """Create a new task on the task board."""
    try:
        from core.task_board import TaskBoard
        board = TaskBoard()
        task = board.create(description)
        return {"ok": True, "task_id": task.task_id, "description": task.description,
                "status": task.status}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_task_status(task_id: str = "", **_) -> dict:
    """Get task status from the task board. If no task_id, list recent tasks."""
    try:
        if not os.path.exists(".task_board.json"):
            return {"ok": True, "tasks": [], "total": 0}

        with open(".task_board.json") as f:
            data = json.load(f)

        if task_id:
            # Find by prefix match
            matches = {tid: t for tid, t in data.items()
                       if tid.startswith(task_id)}
            if not matches:
                return {"ok": False, "error": f"No task matching '{task_id}'"}
            return {"ok": True, "tasks": matches}

        # Return last 10 tasks
        items = list(data.items())[-10:]
        recent = {tid: {"status": t["status"],
                        "agent_id": t.get("agent_id", ""),
                        "description": t["description"][:80]}
                  for tid, t in items}
        return {"ok": True, "tasks": recent, "total": len(data)}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_send_mail(to: str, content: str,
                      agent_id: str = "tool", msg_type: str = "message",
                      **_) -> dict:
    """Send a message to another agent's mailbox."""
    try:
        mailbox_dir = ".mailbox"
        os.makedirs(mailbox_dir, exist_ok=True)
        mailbox_file = os.path.join(mailbox_dir, f"{to}.jsonl")

        entry = {
            "from": agent_id,
            "type": msg_type,
            "content": content,
            "ts": datetime.now(timezone.utc).isoformat(),
        }
        with open(mailbox_file, "a") as f:
            f.write(json.dumps(entry) + "\n")

        return {"ok": True, "to": to, "from": agent_id}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_generate_doc(**kwargs) -> dict:
    """Generate a document file from content.

    Supported formats (8 total):
      - pdf:  Styled PDF with CJK support, tables, code blocks (via fpdf2).
      - docx: Word document with rich formatting, tables (via python-docx).
      - xlsx: Spreadsheet from JSON rows or markdown tables (via openpyxl).
      - pptx: PowerPoint slides split by ## headings (via python-pptx).
      - csv:  CSV from tabular data (JSON, markdown table, TSV).
      - txt:  Plain text with markdown formatting stripped.
      - md:   Markdown file (pass-through).
      - html: Styled HTML with tables, code blocks, lists.

    PDF auto-falls-back to DOCX if generation fails (e.g. font issues).
    The generated file is written to output_path (defaults to /tmp/).
    After generation, use send_file to deliver to the user.
    """
    fmt = (kwargs.get("format") or "pdf").lower().strip()
    content = kwargs.get("content", "")
    output_path = kwargs.get("output_path", "")
    title = kwargs.get("title", "")
    agent_id = kwargs.get("_agent_id", "unknown")

    # If minimax adapter recovered params from truncated JSON, log it
    if kwargs.get("_recovered_from_truncation"):
        logger.info("[generate_doc] Using params recovered from truncated tool args "
                    "(format=%s, content_len=%d)", fmt, len(content))

    # Fallback: if LLM adapter couldn't parse JSON (e.g. MiniMax Unicode issue
    # or truncated arguments), attempt to extract content from raw string.
    if not content and "_raw_args" in kwargs:
        try:
            import re as _re
            raw = kwargs["_raw_args"]

            # Strategy 1: regex for complete "content": "..." value
            m = _re.search(r'"content"\s*:\s*"((?:[^"\\]|\\.)*)"', raw, _re.DOTALL)
            if m:
                content = m.group(1).encode().decode('unicode_escape', errors='replace')

            # Strategy 2: if content regex failed (e.g. truncated string),
            # grab everything after "content": " up to the end of the string.
            if not content:
                m2 = _re.search(r'"content"\s*:\s*"(.*)', raw, _re.DOTALL)
                if m2:
                    raw_content = m2.group(1)
                    # Strip trailing incomplete JSON delimiters
                    raw_content = raw_content.rstrip('} \t\n\r')
                    # Remove trailing incomplete escape sequence
                    if raw_content.endswith('\\'):
                        raw_content = raw_content[:-1]
                    # Remove trailing unmatched quote
                    raw_content = raw_content.rstrip('"')
                    # Unescape JSON string escapes (\\n → \n, etc.)
                    try:
                        content = raw_content.encode().decode('unicode_escape', errors='replace')
                    except Exception:
                        content = raw_content.replace('\\n', '\n').replace('\\t', '\t')
                    if content:
                        logger.info("[generate_doc] Recovered truncated content (%d chars)",
                                    len(content))

            # Also try to recover format/title from raw args
            if not fmt or fmt == "pdf":
                mf = _re.search(r'"format"\s*:\s*"([^"]+)"', raw)
                if mf:
                    fmt = mf.group(1).lower().strip()
            if not title:
                mt = _re.search(r'"title"\s*:\s*"((?:[^"\\]|\\.)*)"', raw)
                if mt:
                    title = mt.group(1)
            if content:
                logger.info("[generate_doc] Recovered params from raw args (parse_error: %s)",
                            kwargs.get("_parse_error", "?"))
        except Exception as exc:
            logger.warning("[generate_doc] Failed to recover from raw args: %s", exc)

    if not content:
        return {"ok": False, "error": "content parameter is required"}

    # Normalise format aliases
    _FMT_ALIASES = {"excel": "xlsx", "word": "docx", "powerpoint": "pptx",
                    "ppt": "pptx", "text": "txt", "markdown": "md",
                    "htm": "html"}
    fmt = _FMT_ALIASES.get(fmt, fmt)

    # Default output path
    if not output_path:
        ts = int(time.time())
        ext = {"xlsx": "xlsx", "docx": "docx", "pptx": "pptx",
               "csv": "csv", "txt": "txt", "md": "md",
               "html": "html", "pdf": "pdf"}.get(fmt, fmt)
        output_path = f"/tmp/doc_{ts}.{ext}"

    _audit_log("generate_doc", agent_id=agent_id, format=fmt, path=output_path)

    try:
        if fmt == "pdf":
            result = _gen_pdf(content, output_path, title)
        elif fmt == "xlsx":
            result = _gen_xlsx(content, output_path, title)
        elif fmt == "docx":
            result = _gen_docx(content, output_path, title)
        elif fmt == "pptx":
            result = _gen_pptx(content, output_path, title)
        elif fmt == "csv":
            result = _gen_csv(content, output_path, title)
        elif fmt == "txt":
            result = _gen_txt(content, output_path, title)
        elif fmt == "md":
            result = _gen_md(content, output_path, title)
        elif fmt == "html":
            result = _gen_html(content, output_path, title)
        else:
            return {"ok": False,
                    "error": f"Unsupported format: {fmt}. "
                             "Use pdf, docx, xlsx, pptx, csv, txt, md, or html."}
    except ImportError as e:
        return {"ok": False,
                "error": f"Missing library for {fmt}: {e}. Install via pip3."}
    except Exception as e:
        # ── Auto-fallback: if PDF fails, retry as DOCX ──
        if fmt == "pdf":
            logger.warning("PDF generation failed (%s), falling back to DOCX", e)
            try:
                fallback_path = output_path.rsplit(".", 1)[0] + ".docx"
                result = _gen_docx(content, fallback_path, title)
                result["original_format"] = "pdf"
                result["fallback"] = True
                result["fallback_reason"] = str(e)
                logger.info("Fallback DOCX generated: %s", fallback_path)
            except Exception as e2:
                logger.exception("generate_doc failed (PDF + DOCX fallback)")
                return {"ok": False,
                        "error": f"PDF failed: {e}; DOCX fallback also failed: {e2}"}
        else:
            logger.exception("generate_doc failed")
            return {"ok": False, "error": str(e)}

    # ── Auto-send: if channel session is active, deliver file immediately ──
    if result.get("ok") and result.get("path"):
        try:
            from adapters.channels.manager import ChannelManager
            session = ChannelManager.get_active_session()
            if session and session.get("session_id"):
                send_result = _handle_send_file(
                    file_path=result["path"],
                    caption=title or f"📄 {os.path.basename(result['path'])}")
                result["auto_sent"] = send_result.get("ok", False)
                if send_result.get("ok"):
                    result["delivery"] = "sent"
                    result["message"] = (
                        f"File generated and sent via {session['channel']}")
                else:
                    result["delivery"] = "failed"
                    result["send_error"] = send_result.get("error", "")
                    result["retry_hint"] = _file_delivery_hint(
                        result["path"], "auto-send failed",
                        send_result.get("error", "unknown"))
                    logger.warning("generate_doc auto-send failed: %s",
                                   send_result.get("error"))
            else:
                result["delivery"] = "no_session"
                result["retry_hint"] = _file_delivery_hint(
                    result["path"], "no active channel session")
        except Exception as e:
            logger.warning("generate_doc auto-send error: %s", e)
            result["delivery"] = "manual"
            result["retry_hint"] = _file_delivery_hint(
                result["path"], "auto-send exception", str(e))

    return result


def _file_delivery_hint(path: str, reason: str, error: str = "") -> str:
    """Build a standard file delivery retry hint message."""
    msg = f"File generated at {path}, but {reason}"
    if error:
        msg += f" (error: {error})"
    msg += f". Call send_file(file_path='{path}') to retry."
    return msg


def _gen_pdf(content: str, output_path: str, title: str) -> dict:
    """Generate PDF from text/markdown content using fpdf2.

    Improvements over v1:
      - Robust CJK font chain (macOS + Linux), no deprecated uni=True
      - All text uses multi_cell() for proper wrapping (prevents
        "Not enough horizontal space" crash on long CJK lines)
      - Markdown table rendering (| col | col |)
      - Bold / italic inline formatting via **text** and *text*
      - Stripped markdown formatting chars (**bold**, *italic*) in plain text
    """
    import re as _re
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.add_page()

    # ── Font setup: CJK-capable Unicode fonts ──
    _FONT_CHAIN = [
        # macOS
        "/Library/Fonts/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/System/Library/Fonts/STHeiti Medium.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
        "/System/Library/Fonts/Supplemental/Songti.ttc",
        # Linux
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    font_name = "Helvetica"
    for fp in _FONT_CHAIN:
        if os.path.exists(fp):
            try:
                pdf.add_font("UniFont", "", fp)
                pdf.set_font("UniFont", size=11)
                font_name = "UniFont"
                break
            except Exception:
                continue
    if font_name == "Helvetica":
        pdf.set_font("Helvetica", size=11)

    def _has_cjk(text: str) -> bool:
        return any('\u4e00' <= ch <= '\u9fff' or '\u3000' <= ch <= '\u30ff'
                   or '\uac00' <= ch <= '\ud7af' for ch in text)

    def _safe_text(text: str) -> str:
        """If using Helvetica (no CJK), replace CJK chars with ?."""
        if font_name != "Helvetica":
            return text
        return "".join(ch if ord(ch) < 0x2E80 else "?" for ch in text)

    def _strip_md_inline(text: str) -> str:
        """Strip **bold** and *italic* markers for plain rendering."""
        text = _re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = _re.sub(r'\*(.+?)\*', r'\1', text)
        text = _re.sub(r'`(.+?)`', r'\1', text)
        return text

    # ── Title ──
    if title:
        pdf.set_font_size(18)
        pdf.multi_cell(0, 12, _safe_text(_strip_md_inline(title)),
                        align="C", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(6)
        pdf.set_font_size(11)

    # ── Pre-process: detect table blocks ──
    lines = content.split("\n")
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()

        # ── Markdown table block ──
        if "|" in stripped and stripped.startswith("|") and stripped.endswith("|"):
            table_lines = []
            while i < len(lines):
                row = lines[i].strip()
                if not ("|" in row and row.startswith("|")):
                    break
                # Skip separator rows (|---|---|)
                cells = [c.strip() for c in row.split("|")[1:-1]]
                if cells and not all(_re.match(r'^[-:]+$', c) for c in cells):
                    table_lines.append(cells)
                i += 1
            if table_lines:
                _render_pdf_table(pdf, table_lines, font_name, _safe_text,
                                  _strip_md_inline)
                pdf.ln(4)
            continue

        # ── Headers ──
        if stripped.startswith("### "):
            pdf.ln(4)
            pdf.set_font_size(13)
            pdf.multi_cell(0, 8, _safe_text(_strip_md_inline(stripped[4:])),
                           new_x="LMARGIN", new_y="NEXT")
            pdf.set_font_size(11)
        elif stripped.startswith("## "):
            pdf.ln(5)
            pdf.set_font_size(15)
            pdf.multi_cell(0, 9, _safe_text(_strip_md_inline(stripped[3:])),
                           new_x="LMARGIN", new_y="NEXT")
            pdf.set_font_size(11)
        elif stripped.startswith("# "):
            pdf.ln(6)
            pdf.set_font_size(17)
            pdf.multi_cell(0, 10, _safe_text(_strip_md_inline(stripped[2:])),
                           new_x="LMARGIN", new_y="NEXT")
            pdf.set_font_size(11)

        # ── Bullet list ──
        elif stripped.startswith(("- ", "* ", "• ")):
            bullet_text = stripped.lstrip("-*• ").strip()
            pdf.set_x(pdf.l_margin + 8)
            pdf.multi_cell(0, 7,
                           _safe_text("•  " + _strip_md_inline(bullet_text)),
                           new_x="LMARGIN", new_y="NEXT")

        # ── Horizontal rule ──
        elif stripped.startswith(("---", "***", "___")):
            pdf.ln(3)
            x = pdf.l_margin
            pdf.line(x, pdf.get_y(), x + 170, pdf.get_y())
            pdf.ln(3)

        # ── Blank line ──
        elif stripped == "":
            pdf.ln(4)

        # ── Code block ──
        elif stripped.startswith("```"):
            # Consume until closing ```
            i += 1
            code_lines = []
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            if code_lines:
                pdf.set_font_size(9)
                for cl in code_lines:
                    pdf.multi_cell(0, 5, _safe_text(cl),
                                   new_x="LMARGIN", new_y="NEXT")
                pdf.set_font_size(11)
                pdf.ln(2)

        else:
            # ── Numbered list ──
            num_match = _re.match(r'^(\d+)[.)]\s+(.+)$', stripped)
            if num_match:
                pdf.multi_cell(
                    0, 7,
                    _safe_text(f"{num_match.group(1)}. "
                               f"{_strip_md_inline(num_match.group(2))}"),
                    new_x="LMARGIN", new_y="NEXT")
            else:
                pdf.multi_cell(0, 7, _safe_text(_strip_md_inline(stripped)),
                               new_x="LMARGIN", new_y="NEXT")

        i += 1

    os.makedirs(os.path.dirname(output_path) or "/tmp", exist_ok=True)
    pdf.output(output_path)
    size = os.path.getsize(output_path)
    return {"ok": True, "path": output_path, "format": "pdf",
            "size": size, "pages": pdf.pages_count}


def _render_pdf_table(pdf, rows, font_name, _safe_text, _strip_md_inline):
    """Render a markdown table into the PDF using fpdf2 built-in table API."""
    if not rows:
        return
    n_cols = max(len(r) for r in rows)
    with pdf.table(first_row_as_headings=False,
                   col_widths=tuple(1 for _ in range(n_cols))) as table:
        for row_data in rows:
            row = table.row()
            for c_idx in range(n_cols):
                text = _safe_text(_strip_md_inline(
                    str(row_data[c_idx]) if c_idx < len(row_data) else ""))
                row.cell(text)


def _gen_xlsx(content: str, output_path: str, title: str) -> dict:
    """Generate Excel spreadsheet from JSON content using openpyxl."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment

    wb = Workbook()
    ws = wb.active
    ws.title = title or "Sheet1"

    # Parse content: accept JSON array of arrays, or line-based tabular data
    rows = []
    try:
        parsed = json.loads(content)
        if isinstance(parsed, list):
            rows = parsed
        elif isinstance(parsed, dict):
            # {"headers": [...], "rows": [[...], ...]}
            if "headers" in parsed:
                rows = [parsed["headers"]] + parsed.get("rows", [])
            else:
                # Single dict → one row
                rows = [list(parsed.keys()), list(parsed.values())]
    except (json.JSONDecodeError, TypeError):
        # Fallback: parse as TSV/CSV-like lines
        for line in content.strip().split("\n"):
            if "|" in line:
                cells = [c.strip() for c in line.split("|") if c.strip()]
                if cells and not all(c.startswith("-") for c in cells):
                    rows.append(cells)
            elif "\t" in line:
                rows.append(line.split("\t"))
            elif "," in line:
                rows.append(line.split(","))
            else:
                rows.append([line])

    if not rows:
        return {"ok": False, "error": "No tabular data found in content"}

    # Write rows
    for r_idx, row in enumerate(rows):
        if not isinstance(row, (list, tuple)):
            row = [row]
        for c_idx, val in enumerate(row):
            cell = ws.cell(row=r_idx + 1, column=c_idx + 1, value=val)
            if r_idx == 0:  # Header styling
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill(start_color="4472C4",
                                        end_color="4472C4",
                                        fill_type="solid")
                cell.alignment = Alignment(horizontal="center")

    # Auto-width columns
    for col in ws.columns:
        max_len = max((len(str(cell.value or "")) for cell in col), default=8)
        ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)

    os.makedirs(os.path.dirname(output_path) or "/tmp", exist_ok=True)
    wb.save(output_path)
    size = os.path.getsize(output_path)
    return {"ok": True, "path": output_path, "format": "xlsx",
            "size": size, "rows": len(rows)}


def _gen_docx(content: str, output_path: str, title: str) -> dict:
    """Generate Word document from text/markdown content using python-docx.

    Improvements:
      - CJK-friendly font (PingFang SC / SimSun / Noto Sans CJK)
      - Markdown table rendering (| col | col |)
      - Bold / italic inline formatting
      - Code block support
    """
    import re as _re
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.oxml.ns import qn

    doc = Document()

    # ── CJK-friendly default font ──
    style = doc.styles["Normal"]
    font = style.font
    font.size = Pt(11)
    font.name = "Arial Unicode MS"
    # Set East Asian font for CJK rendering
    rpr = style.element.get_or_add_rPr()
    ea_font = rpr.makeelement(qn('w:rFonts'), {
        qn('w:eastAsia'): 'PingFang SC',
    })
    rpr.insert(0, ea_font)

    def _strip_md(text):
        text = _re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = _re.sub(r'\*(.+?)\*', r'\1', text)
        text = _re.sub(r'`(.+?)`', r'\1', text)
        return text

    def _add_rich_paragraph(doc, text, style_name=None):
        """Add paragraph with **bold** and *italic* inline formatting."""
        p = doc.add_paragraph(style=style_name)
        parts = _re.split(r'(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`)', text)
        for part in parts:
            if part.startswith("**") and part.endswith("**"):
                run = p.add_run(part[2:-2])
                run.bold = True
            elif part.startswith("*") and part.endswith("*"):
                run = p.add_run(part[1:-1])
                run.italic = True
            elif part.startswith("`") and part.endswith("`"):
                run = p.add_run(part[1:-1])
                run.font.name = "Courier New"
                run.font.size = Pt(10)
            else:
                p.add_run(part)
        return p

    if title:
        doc.add_heading(_strip_md(title), level=0)

    lines = content.split("\n")
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()

        # ── Markdown table ──
        if ("|" in stripped and stripped.startswith("|")
                and stripped.endswith("|")):
            table_rows = []
            while i < len(lines):
                row = lines[i].strip()
                if not ("|" in row and row.startswith("|")):
                    break
                cells = [c.strip() for c in row.split("|")[1:-1]]
                if cells and not all(_re.match(r'^[-:]+$', c) for c in cells):
                    table_rows.append(cells)
                i += 1
            if table_rows:
                n_cols = max(len(r) for r in table_rows)
                tbl = doc.add_table(rows=len(table_rows), cols=n_cols,
                                    style="Light Grid Accent 1")
                for r_idx, row_data in enumerate(table_rows):
                    for c_idx, cell_text in enumerate(row_data):
                        if c_idx < n_cols:
                            tbl.cell(r_idx, c_idx).text = _strip_md(cell_text)
                            # Bold header row
                            if r_idx == 0:
                                for run in tbl.cell(r_idx, c_idx).paragraphs[0].runs:
                                    run.bold = True
                doc.add_paragraph("")
            continue

        # ── Headers ──
        if stripped.startswith("### "):
            doc.add_heading(_strip_md(stripped[4:]), level=3)
        elif stripped.startswith("## "):
            doc.add_heading(_strip_md(stripped[3:]), level=2)
        elif stripped.startswith("# "):
            doc.add_heading(_strip_md(stripped[2:]), level=1)

        # ── Bullet list ──
        elif stripped.startswith(("- ", "* ", "• ")):
            text = stripped.lstrip("-*• ").strip()
            _add_rich_paragraph(doc, text, style_name="List Bullet")

        # ── Horizontal rule ──
        elif stripped.startswith(("---", "***", "___")):
            doc.add_paragraph("_" * 50)

        # ── Code block ──
        elif stripped.startswith("```"):
            i += 1
            code_lines = []
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            if code_lines:
                p = doc.add_paragraph()
                run = p.add_run("\n".join(code_lines))
                run.font.name = "Courier New"
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # ── Blank line ──
        elif stripped == "":
            doc.add_paragraph("")

        else:
            # ── Numbered list ──
            num_match = _re.match(r'^(\d+)[.)]\s+(.+)$', stripped)
            if num_match:
                _add_rich_paragraph(doc, num_match.group(2),
                                    style_name="List Number")
            else:
                _add_rich_paragraph(doc, stripped)

        i += 1

    os.makedirs(os.path.dirname(output_path) or "/tmp", exist_ok=True)
    doc.save(output_path)
    size = os.path.getsize(output_path)
    return {"ok": True, "path": output_path, "format": "docx", "size": size}


def _gen_pptx(content: str, output_path: str, title: str) -> dict:
    """Generate PowerPoint presentation from text/markdown content.

    Splits content by ## headings into slides. Each slide gets a title
    and bullet-point body. Supports CJK text natively.
    """
    import re as _re
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Title slide ──
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = _re.sub(r'\*\*(.+?)\*\*', r'\1', title or "Untitled")
    if slide.placeholders[1]:
        slide.placeholders[1].text = ""

    # ── Split content into slide sections by ## headings ──
    sections: list[dict] = []
    current: dict | None = None
    for line in content.split("\n"):
        stripped = line.strip()
        if stripped.startswith("## "):
            if current:
                sections.append(current)
            current = {"title": _re.sub(r'\*\*(.+?)\*\*', r'\1', stripped[3:]),
                        "body": []}
        elif stripped.startswith("# ") and not current:
            # Top-level heading as subtitle on title slide
            if slide.placeholders[1]:
                slide.placeholders[1].text = _re.sub(
                    r'\*\*(.+?)\*\*', r'\1', stripped[2:])
        elif current is not None:
            if stripped and not stripped.startswith(("---", "***", "___", "```")):
                # Clean markdown formatting
                clean = _re.sub(r'\*\*(.+?)\*\*', r'\1', stripped)
                clean = _re.sub(r'\*(.+?)\*', r'\1', clean)
                clean = _re.sub(r'`(.+?)`', r'\1', clean)
                clean = clean.lstrip("-*• ").strip()
                if clean:
                    current["body"].append(clean)
    if current:
        sections.append(current)

    # If no ## headings found, create one slide with all content
    if not sections:
        body_lines = []
        for line in content.split("\n"):
            stripped = line.strip()
            if stripped and not stripped.startswith(("#", "---", "***", "```")):
                clean = _re.sub(r'\*\*(.+?)\*\*', r'\1', stripped)
                clean = clean.lstrip("-*• ").strip()
                if clean:
                    body_lines.append(clean)
        if body_lines:
            sections = [{"title": title or "Content", "body": body_lines}]

    # ── Render slides ──
    slide_layout = prs.slide_layouts[1]  # Title and Content
    for sec in sections:
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = sec["title"]
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for idx, bullet in enumerate(sec["body"]):
            if idx == 0:
                tf.paragraphs[0].text = bullet
                tf.paragraphs[0].font.size = Pt(18)
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)

    # ── Markdown table → table slide ──
    table_rows = []
    in_table = False
    for line in content.split("\n"):
        stripped = line.strip()
        if "|" in stripped and stripped.startswith("|") and stripped.endswith("|"):
            cells = [c.strip() for c in stripped.split("|")[1:-1]]
            if cells and not all(_re.match(r'^[-:]+$', c) for c in cells):
                table_rows.append(cells)
                in_table = True
        elif in_table:
            break
    if table_rows and len(table_rows) > 1:
        slide_layout = prs.slide_layouts[5]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        n_cols = max(len(r) for r in table_rows)
        tbl_shape = slide.shapes.add_table(
            len(table_rows), n_cols,
            Inches(0.5), Inches(0.5),
            Inches(12), Inches(6))
        tbl = tbl_shape.table
        for r_idx, row_data in enumerate(table_rows):
            for c_idx in range(n_cols):
                cell_text = row_data[c_idx] if c_idx < len(row_data) else ""
                tbl.cell(r_idx, c_idx).text = _re.sub(
                    r'\*\*(.+?)\*\*', r'\1', cell_text)

    os.makedirs(os.path.dirname(output_path) or "/tmp", exist_ok=True)
    prs.save(output_path)
    size = os.path.getsize(output_path)
    return {"ok": True, "path": output_path, "format": "pptx",
            "size": size, "slides": len(prs.slides)}


def _gen_csv(content: str, output_path: str, title: str) -> dict:
    """Generate CSV from tabular content (markdown table, JSON, or TSV)."""
    import csv
    import io
    import re as _re

    rows: list[list[str]] = []

    # Try JSON first
    try:
        parsed = json.loads(content)
        if isinstance(parsed, list):
            for item in parsed:
                if isinstance(item, (list, tuple)):
                    rows.append([str(v) for v in item])
                elif isinstance(item, dict):
                    if not rows:
                        rows.append(list(item.keys()))
                    rows.append([str(v) for v in item.values()])
    except (json.JSONDecodeError, TypeError):
        pass

    # Fallback: parse markdown table or TSV
    if not rows:
        for line in content.strip().split("\n"):
            stripped = line.strip()
            if "|" in stripped:
                cells = [c.strip() for c in stripped.split("|")]
                # Remove empty leading/trailing from |col|col|
                if cells and cells[0] == "":
                    cells = cells[1:]
                if cells and cells[-1] == "":
                    cells = cells[:-1]
                if cells and not all(_re.match(r'^[-:]+$', c) for c in cells):
                    rows.append(cells)
            elif "\t" in stripped:
                rows.append(stripped.split("\t"))
            elif stripped:
                rows.append([stripped])

    if not rows:
        return {"ok": False, "error": "No tabular data found in content"}

    os.makedirs(os.path.dirname(output_path) or "/tmp", exist_ok=True)
    with open(output_path, "w", newline="", encoding="utf-8-sig") as f:
        writer = csv.writer(f)
        writer.writerows(rows)

    size = os.path.getsize(output_path)
    return {"ok": True, "path": output_path, "format": "csv",
            "size": size, "rows": len(rows)}


def _gen_txt(content: str, output_path: str, title: str) -> dict:
    """Generate plain text file, stripping markdown formatting."""
    import re as _re

    lines = []
    if title:
        lines.append(title)
        lines.append("=" * len(title))
        lines.append("")

    for line in content.split("\n"):
        # Strip markdown formatting but keep structure
        clean = line
        clean = _re.sub(r'^#{1,6}\s+', '', clean)          # Headers
        clean = _re.sub(r'\*\*(.+?)\*\*', r'\1', clean)    # Bold
        clean = _re.sub(r'\*(.+?)\*', r'\1', clean)        # Italic
        clean = _re.sub(r'`(.+?)`', r'\1', clean)          # Inline code
        lines.append(clean)

    os.makedirs(os.path.dirname(output_path) or "/tmp", exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    size = os.path.getsize(output_path)
    return {"ok": True, "path": output_path, "format": "txt", "size": size}


def _gen_md(content: str, output_path: str, title: str) -> dict:
    """Generate Markdown file (pass-through with optional title)."""
    text = ""
    if title:
        text = f"# {title}\n\n"
    text += content

    os.makedirs(os.path.dirname(output_path) or "/tmp", exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)

    size = os.path.getsize(output_path)
    return {"ok": True, "path": output_path, "format": "md", "size": size}


def _gen_html(content: str, output_path: str, title: str) -> dict:
    """Generate styled HTML file from markdown-like content."""
    import re as _re
    import html as _html

    html_lines = [
        '<!DOCTYPE html>',
        '<html lang="zh"><head><meta charset="utf-8">',
        f'<title>{_html.escape(title or "Document")}</title>',
        '<style>',
        'body{font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,'
        '"Helvetica Neue",Arial,sans-serif;max-width:800px;margin:40px auto;'
        'padding:0 20px;line-height:1.6;color:#1d1d1f}',
        'h1,h2,h3{margin-top:1.5em}',
        'table{border-collapse:collapse;width:100%;margin:1em 0}',
        'th,td{border:1px solid #ddd;padding:8px 12px;text-align:left}',
        'th{background:#f5f5f7;font-weight:600}',
        'code{background:#f0f0f0;padding:2px 6px;border-radius:4px;font-size:0.9em}',
        'pre{background:#f5f5f7;padding:16px;border-radius:8px;overflow-x:auto}',
        'pre code{background:none;padding:0}',
        'hr{border:none;border-top:1px solid #ddd;margin:2em 0}',
        'ul,ol{padding-left:2em}',
        '</style></head><body>',
    ]

    if title:
        html_lines.append(f'<h1>{_html.escape(title)}</h1>')

    lines = content.split("\n")
    i = 0
    in_list = False
    list_type = None

    while i < len(lines):
        stripped = lines[i].strip()

        # ── Table block ──
        if "|" in stripped and stripped.startswith("|") and stripped.endswith("|"):
            html_lines.append('<table>')
            first_row = True
            while i < len(lines):
                row = lines[i].strip()
                if not ("|" in row and row.startswith("|")):
                    break
                cells = [c.strip() for c in row.split("|")[1:-1]]
                if cells and all(_re.match(r'^[-:]+$', c) for c in cells):
                    i += 1
                    continue
                tag = "th" if first_row else "td"
                html_lines.append('<tr>' + ''.join(
                    f'<{tag}>{_html.escape(c)}</{tag}>' for c in cells) + '</tr>')
                first_row = False
                i += 1
            html_lines.append('</table>')
            continue

        # Close open list if needed
        if in_list and not stripped.startswith(("- ", "* ", "• ")) and \
                not _re.match(r'^\d+[.)]\s', stripped):
            html_lines.append(f'</{list_type}>')
            in_list = False

        # ── Code block ──
        if stripped.startswith("```"):
            i += 1
            code = []
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code.append(_html.escape(lines[i]))
                i += 1
            html_lines.append('<pre><code>' + '\n'.join(code) + '</code></pre>')
            i += 1
            continue

        # ── Headers ──
        if stripped.startswith("### "):
            html_lines.append(f'<h3>{_html.escape(stripped[4:])}</h3>')
        elif stripped.startswith("## "):
            html_lines.append(f'<h2>{_html.escape(stripped[3:])}</h2>')
        elif stripped.startswith("# "):
            html_lines.append(f'<h1>{_html.escape(stripped[2:])}</h1>')
        # ── HR ──
        elif stripped.startswith(("---", "***", "___")):
            html_lines.append('<hr>')
        # ── Bullet list ──
        elif stripped.startswith(("- ", "* ", "• ")):
            if not in_list or list_type != "ul":
                if in_list:
                    html_lines.append(f'</{list_type}>')
                html_lines.append('<ul>')
                in_list = True
                list_type = "ul"
            text = _html.escape(stripped.lstrip("-*• ").strip())
            html_lines.append(f'<li>{text}</li>')
        # ── Numbered list ──
        elif _re.match(r'^\d+[.)]\s', stripped):
            if not in_list or list_type != "ol":
                if in_list:
                    html_lines.append(f'</{list_type}>')
                html_lines.append('<ol>')
                in_list = True
                list_type = "ol"
            text = _re.sub(r'^\d+[.)]\s+', '', stripped)
            html_lines.append(f'<li>{_html.escape(text)}</li>')
        # ── Blank ──
        elif stripped == "":
            html_lines.append('')
        # ── Paragraph ──
        else:
            html_lines.append(f'<p>{_html.escape(stripped)}</p>')

        i += 1

    if in_list:
        html_lines.append(f'</{list_type}>')

    html_lines.append('</body></html>')

    os.makedirs(os.path.dirname(output_path) or "/tmp", exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write('\n'.join(html_lines))

    size = os.path.getsize(output_path)
    return {"ok": True, "path": output_path, "format": "html", "size": size}


def _handle_send_file(**kwargs) -> dict:
    """Send a file to the user via their chat channel.

    Routes through the gateway's /v1/send_file proxy which relays to
    the ChannelManager (running in the gateway process).  Falls back
    to the .file_delivery/ queue if the gateway is unreachable.
    """
    file_path = kwargs.get("file_path", "")
    caption = kwargs.get("caption", "")

    if not file_path:
        return {"ok": False, "error": "file_path parameter required"}

    abs_path = os.path.abspath(file_path)
    if not os.path.exists(abs_path):
        return {"ok": False, "error": f"File not found: {file_path}"}

    # Get active channel session
    try:
        from adapters.channels.manager import ChannelManager
        session = ChannelManager.get_active_session()
    except ImportError:
        session = None
    if not session:
        # Fallback: try reading the file directly
        session_path = ".channel_session.json"
        if os.path.exists(session_path):
            try:
                with open(session_path, "r") as f:
                    session = json.load(f)
            except Exception:
                pass

    # Session staleness check
    if session:
        session_age = time.time() - session.get("ts", 0)
        if session_age > 3600:
            logger.warning("send_file: channel session is %.0fm old, may be stale. "
                           "User should send a new message to refresh.",
                           session_age / 60)

    if not session or not session.get("session_id"):
        return {"ok": False,
                "error": "No active channel session. "
                "Fix: (1) run `cleo gateway start`, "
                "(2) send a message from Telegram/Discord to establish session, "
                "(3) retry. File saved at: " + abs_path}

    session_id = session["session_id"]

    # ── Primary path: HTTP proxy to gateway ──
    gateway_port = int(os.environ.get("CLEO_GATEWAY_PORT", "19789"))
    gateway_token = os.environ.get("CLEO_GATEWAY_TOKEN", "")
    # Also try reading token from file
    if not gateway_token:
        try:
            token_path = ".gateway_token"
            if os.path.exists(token_path):
                with open(token_path) as f:
                    gateway_token = f.read().strip()
        except Exception:
            pass

    import urllib.request

    # Gateway health pre-check
    try:
        health_req = urllib.request.Request(
            f"http://127.0.0.1:{gateway_port}/health", method="GET")
        urllib.request.urlopen(health_req, timeout=2)
    except Exception as e:
        logger.warning("send_file: gateway unreachable at port %d (%s), "
                       "will fall back to queue", gateway_port, e)

    try:
        payload = json.dumps({
            "session_id": session_id,
            "file_path": abs_path,
            "caption": caption,
        }).encode()
        req = urllib.request.Request(
            f"http://127.0.0.1:{gateway_port}/v1/send_file",
            data=payload,
            headers={
                "Content-Type": "application/json",
                "Authorization": f"Bearer {gateway_token}" if gateway_token else "",
            },
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=30) as resp:
            result = json.loads(resp.read())
        if result.get("ok"):
            return {"ok": True, "message": f"File sent via {session_id}",
                    "message_id": result.get("message_id", "")}
        else:
            return {"ok": False, "error": result.get("error", "Gateway returned error")}
    except Exception as e:
        logger.warning("send_file HTTP proxy failed (%s), trying direct path", e)

    # ── Fallback 1: Direct ChannelManager call (if running in same process) ──
    try:
        from adapters.channels.manager import ChannelManager
        cm = getattr(ChannelManager, '_instance', None)
        if cm is None:
            # Try the module-level global in gateway
            import sys
            gw_mod = sys.modules.get("core.gateway")
            if gw_mod:
                cm = getattr(gw_mod, '_channel_manager', None)
        if cm and hasattr(cm, '_loop') and cm._loop and cm._loop.is_running():
            import asyncio as _asyncio
            future = _asyncio.run_coroutine_threadsafe(
                cm.send_file(session_id, abs_path, caption),
                cm._loop)
            msg_id = future.result(timeout=30)
            if msg_id:
                return {"ok": True, "message": f"File sent directly via {session_id}",
                        "message_id": msg_id, "method": "direct"}
            else:
                logger.warning("send_file direct path: send_file returned empty msg_id")
    except Exception as direct_err:
        logger.debug("send_file direct path failed: %s", direct_err)

    # ── Fallback 2: write to .file_delivery/ queue (consumed by gateway poller) ──
    try:
        delivery_dir = ".file_delivery"
        os.makedirs(delivery_dir, exist_ok=True)
        delivery = {
            "file_path": abs_path,
            "caption": caption,
            "session_id": session_id,
            "channel": session.get("channel", ""),
            "chat_id": session.get("chat_id", ""),
            "ts": time.time(),
            "retry_count": 0,
        }
        delivery_file = os.path.join(delivery_dir, f"{int(time.time()*1000)}.json")
        with open(delivery_file, "w") as f:
            json.dump(delivery, f)
        return {"ok": True, "message": f"File queued for delivery: {file_path}",
                "delivery_id": os.path.basename(delivery_file)}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_message(**kwargs) -> dict:
    """Send a message (text, file, or both) to the user via their chat channel.

    Unified messaging tool inspired by OpenClaw's message() pattern.
    Three delivery tiers:
      1. HTTP proxy → /v1/send_message  (text) + /v1/send_file (file)
      2. Direct ChannelManager call
      3. .file_delivery/ queue fallback
    """
    text = kwargs.get("text", "")
    file_path = kwargs.get("file_path", "")
    caption = kwargs.get("caption", "")

    if not text and not file_path:
        return {"ok": False, "error": "At least one of 'text' or 'file_path' is required"}

    # Validate file exists (if provided)
    abs_file_path = ""
    if file_path:
        abs_file_path = os.path.abspath(file_path)
        if not os.path.exists(abs_file_path):
            return {"ok": False, "error": f"File not found: {file_path}"}

    # ── Resolve active channel session (same as send_file) ──
    try:
        from adapters.channels.manager import ChannelManager
        session = ChannelManager.get_active_session()
    except ImportError:
        session = None
    if not session:
        session_path = ".channel_session.json"
        if os.path.exists(session_path):
            try:
                with open(session_path, "r") as f:
                    session = json.load(f)
            except Exception:
                pass

    if session:
        session_age = time.time() - session.get("ts", 0)
        if session_age > 3600:
            logger.warning("message: channel session is %.0fm old, may be stale.",
                           session_age / 60)

    if not session or not session.get("session_id"):
        return {"ok": False,
                "error": "No active channel session. "
                "Fix: (1) run `cleo gateway start`, "
                "(2) send a message from Telegram/Discord to establish session, "
                "(3) retry."}

    session_id = session["session_id"]
    results = {}

    # ── Step 1: Send text (if provided) ──
    if text:
        text_result = _send_text_message(session_id, text, session)
        results["text"] = text_result
        if not text_result.get("ok"):
            logger.warning("message: text send failed: %s", text_result.get("error"))

    # ── Step 2: Send file (if provided) — delegate to _handle_send_file ──
    if file_path:
        file_result = _handle_send_file(file_path=file_path, caption=caption)
        results["file"] = file_result
        if not file_result.get("ok"):
            logger.warning("message: file send failed: %s", file_result.get("error"))

    # Summarize
    text_ok = results.get("text", {}).get("ok", True)  # True if text not sent
    file_ok = results.get("file", {}).get("ok", True)   # True if file not sent
    if text_ok and file_ok:
        parts = []
        if text:
            parts.append("text")
        if file_path:
            parts.append("file")
        return {"ok": True,
                "message": f"Sent {' + '.join(parts)} via {session_id}",
                "details": results}
    else:
        return {"ok": False,
                "error": "Partial failure",
                "details": results}


def _send_text_message(session_id: str, text: str, session: dict) -> dict:
    """Send text via 3-tier delivery: HTTP proxy → direct → queue."""
    gateway_port = int(os.environ.get("CLEO_GATEWAY_PORT", "19789"))
    gateway_token = os.environ.get("CLEO_GATEWAY_TOKEN", "")
    if not gateway_token:
        try:
            token_path = ".gateway_token"
            if os.path.exists(token_path):
                with open(token_path) as f:
                    gateway_token = f.read().strip()
        except Exception:
            pass

    import urllib.request

    # ── Tier 1: HTTP proxy to gateway ──
    try:
        payload = json.dumps({
            "session_id": session_id,
            "text": text,
        }).encode()
        req = urllib.request.Request(
            f"http://127.0.0.1:{gateway_port}/v1/send_message",
            data=payload,
            headers={
                "Content-Type": "application/json",
                "Authorization": f"Bearer {gateway_token}" if gateway_token else "",
            },
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=30) as resp:
            result = json.loads(resp.read())
        if result.get("ok"):
            return {"ok": True, "message": f"Text sent via {session_id}",
                    "message_id": result.get("message_id", "")}
        else:
            return {"ok": False, "error": result.get("error", "Gateway returned error")}
    except Exception as e:
        logger.warning("message text HTTP proxy failed (%s), trying direct path", e)

    # ── Tier 2: Direct ChannelManager call ──
    try:
        from adapters.channels.manager import ChannelManager
        cm = getattr(ChannelManager, '_instance', None)
        if cm is None:
            import sys
            gw_mod = sys.modules.get("core.gateway")
            if gw_mod:
                cm = getattr(gw_mod, '_channel_manager', None)
        if cm and hasattr(cm, '_loop') and cm._loop and cm._loop.is_running():
            import asyncio as _asyncio
            future = _asyncio.run_coroutine_threadsafe(
                cm.send_message(session_id, text),
                cm._loop)
            msg_id = future.result(timeout=30)
            if msg_id:
                return {"ok": True, "message": f"Text sent directly via {session_id}",
                        "message_id": msg_id, "method": "direct"}
    except Exception as direct_err:
        logger.debug("message text direct path failed: %s", direct_err)

    # ── Tier 3: .file_delivery/ queue (now supports text) ──
    try:
        delivery_dir = ".file_delivery"
        os.makedirs(delivery_dir, exist_ok=True)
        delivery = {
            "text": text,
            "session_id": session_id,
            "channel": session.get("channel", ""),
            "chat_id": session.get("chat_id", ""),
            "ts": time.time(),
            "retry_count": 0,
        }
        delivery_file = os.path.join(delivery_dir, f"{int(time.time()*1000)}.json")
        with open(delivery_file, "w") as f:
            json.dump(delivery, f)
        return {"ok": True, "message": "Text queued for delivery",
                "delivery_id": os.path.basename(delivery_file)}
    except Exception as e:
        return {"ok": False, "error": str(e)}


# ── Skill CLI install/check handlers ──

def _handle_check_skill_deps(**kwargs) -> dict:
    """Check which skill CLI dependencies are missing."""
    try:
        from core.skill_deps import (
            scan_skill_deps, get_missing_deps, get_installed_deps,
            check_prerequisites,
        )
    except ImportError:
        return {"ok": False, "error": "skill_deps module not available"}

    skill_name = kwargs.get("skill", "")

    if skill_name:
        # Check a specific skill
        all_deps = scan_skill_deps()
        match = [d for d in all_deps
                 if d["skill"] == skill_name or d["file"] == f"{skill_name}.md"]
        if not match:
            return {"ok": False, "error": f"Skill '{skill_name}' not found or has no CLI deps"}
        dep = match[0]
        return {
            "ok": True,
            "skill": dep["skill"],
            "requires": dep["requires_bins"] or dep.get("requires_any_bins", []),
            "missing": dep["missing_bins"],
            "satisfied": not dep["missing_bins"] and dep["has_any_bin"],
            "install_options": [
                {"kind": e.get("kind"), "label": e.get("label")}
                for e in dep.get("install", [])
            ],
        }

    # Check all skills
    prereqs = check_prerequisites()
    installed = get_installed_deps()
    missing = get_missing_deps()
    return {
        "ok": True,
        "package_managers": prereqs,
        "total": len(installed) + len(missing),
        "installed_count": len(installed),
        "missing_count": len(missing),
        "installed": [
            {"skill": d["skill"], "bins": d["requires_bins"]}
            for d in installed
        ],
        "missing": [
            {"skill": d["skill"],
             "needs": d["missing_bins"] or d.get("requires_any_bins", []),
             "install": [e.get("label") for e in d.get("install", [])]}
            for d in missing
        ],
    }


def _handle_install_skill_cli(**kwargs) -> dict:
    """Install CLI dependencies for a skill. Agents can self-install."""
    try:
        from core.skill_deps import (
            scan_skill_deps, pick_best_installer, install_dep,
            build_install_command,
        )
    except ImportError:
        return {"ok": False, "error": "skill_deps module not available"}

    skill_name = kwargs.get("skill", "")
    if not skill_name:
        return {"ok": False, "error": "skill parameter required"}

    all_deps = scan_skill_deps()
    match = [d for d in all_deps
             if d["skill"] == skill_name or d["file"] == f"{skill_name}.md"]
    if not match:
        return {"ok": False, "error": f"Skill '{skill_name}' not found or has no install entries"}

    dep = match[0]
    if not dep["missing_bins"] and dep["has_any_bin"]:
        return {"ok": True, "message": f"All deps for '{skill_name}' already installed",
                "bins": dep["requires_bins"]}

    install_entries = dep.get("install", [])
    if not install_entries:
        return {"ok": False, "error": f"No install entries for '{skill_name}'"}

    best = pick_best_installer(install_entries)
    if not best:
        return {"ok": False, "error": "No compatible package manager found"}

    cmd = build_install_command(best)
    logger.info("install_skill_cli: %s → %s", skill_name, cmd)

    success = install_dep(best, quiet=False)
    if not success:
        return {"ok": False, "error": f"Install failed: {cmd}",
                "command": cmd}

    # Auto-approve the installed binaries in exec_approvals so agents can use them
    installed_bins = best.get("bins", dep.get("requires_bins", []))
    _auto_approve_skill_bins(installed_bins)

    return {
        "ok": True,
        "skill": skill_name,
        "command": cmd,
        "bins": installed_bins,
        "message": f"Installed {skill_name} via: {cmd}",
    }


def _auto_approve_skill_bins(bins: list[str]):
    """Add installed skill binaries to exec_approvals.json so agents can use them."""
    try:
        from core.exec_tool import add_approval
        for b in bins:
            # Approve the binary (with any args)
            pattern = rf"^{re.escape(b)}\b"
            add_approval(pattern)
            logger.info("Auto-approved exec pattern: %s", pattern)
    except Exception as e:
        logger.warning("Could not auto-approve bins: %s", e)


# ── Remote skill registry handlers ──

def _handle_search_skills(**kwargs) -> dict:
    """Search the remote skill registry for new skills to install."""
    try:
        from core.skill_registry import get_registry
    except ImportError:
        return {"ok": False, "error": "skill_registry module not available"}

    query = kwargs.get("query", "")
    if not query:
        return {"ok": False, "error": "query parameter required"}

    registry = get_registry()
    try:
        results = registry.search(query, limit=kwargs.get("limit", 10))
    except Exception as e:
        return {"ok": False, "error": f"Search failed: {e}"}

    if not results:
        return {
            "ok": True,
            "count": 0,
            "message": f"No skills found matching '{query}'",
            "results": [],
        }

    # Format results for agent consumption
    formatted = []
    for r in results:
        entry = {
            "slug": r.get("slug", ""),
            "name": r.get("name", ""),
            "description": r.get("description", ""),
            "version": r.get("version", ""),
            "tags": r.get("tags", []),
            "installed": r.get("installed", False),
        }
        if r.get("installed"):
            entry["installed_version"] = r.get("installed_version", "")
        if r.get("requires", {}).get("bins"):
            entry["requires_cli"] = r["requires"]["bins"]
        formatted.append(entry)

    return {
        "ok": True,
        "count": len(formatted),
        "query": query,
        "results": formatted,
    }


def _handle_install_remote_skill(**kwargs) -> dict:
    """Install a skill from the remote registry."""
    try:
        from core.skill_registry import get_registry
    except ImportError:
        return {"ok": False, "error": "skill_registry module not available"}

    slug = kwargs.get("slug", "")
    if not slug:
        return {"ok": False, "error": "slug parameter required"}

    agent_id = kwargs.get("agent", "")
    add_to_all = kwargs.get("add_to_all", True)

    registry = get_registry()

    # Step 1: Install the skill
    result = registry.install(slug)
    if not result.get("ok"):
        return result

    # Step 2: Add to agent config
    if add_to_all:
        cfg_result = registry.add_to_all_agents(slug)
    elif agent_id:
        cfg_result = registry.add_to_agent(slug, agent_id)
    else:
        cfg_result = {"ok": True, "message": "Skill installed but not added to any agent config"}

    result["config"] = cfg_result.get("message", "")
    result["message"] = (
        f"{result.get('message', '')}. "
        f"{cfg_result.get('message', '')}. "
        "Skill will be active on the next task (hot-reload)."
    )

    return result


# ── Browser automation handlers ──

def _handle_browser_navigate(**kwargs) -> dict:
    """Navigate the browser to a URL."""
    try:
        from adapters.browser.playwright_adapter import handle_browser_navigate
        return handle_browser_navigate(**kwargs)
    except ImportError:
        return {"ok": False, "error": "playwright not installed. Run: pip install playwright && playwright install chromium"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_browser_click(**kwargs) -> dict:
    """Click an element in the browser."""
    try:
        from adapters.browser.playwright_adapter import handle_browser_click
        return handle_browser_click(**kwargs)
    except ImportError:
        return {"ok": False, "error": "playwright not installed"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_browser_fill(**kwargs) -> dict:
    """Fill a form field in the browser."""
    try:
        from adapters.browser.playwright_adapter import handle_browser_fill
        return handle_browser_fill(**kwargs)
    except ImportError:
        return {"ok": False, "error": "playwright not installed"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_browser_get_text(**kwargs) -> dict:
    """Extract text content from the page."""
    try:
        from adapters.browser.playwright_adapter import handle_browser_get_text
        return handle_browser_get_text(**kwargs)
    except ImportError:
        return {"ok": False, "error": "playwright not installed"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_browser_screenshot(**kwargs) -> dict:
    """Take a screenshot of the page."""
    try:
        from adapters.browser.playwright_adapter import handle_browser_screenshot
        return handle_browser_screenshot(**kwargs)
    except ImportError:
        return {"ok": False, "error": "playwright not installed"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_browser_evaluate(**kwargs) -> dict:
    """Run JavaScript in the page."""
    try:
        from adapters.browser.playwright_adapter import handle_browser_evaluate
        return handle_browser_evaluate(**kwargs)
    except ImportError:
        return {"ok": False, "error": "playwright not installed"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _handle_browser_page_info(**kwargs) -> dict:
    """Get current page info (URL, title)."""
    try:
        from adapters.browser.playwright_adapter import handle_browser_page_info
        return handle_browser_page_info(**kwargs)
    except ImportError:
        return {"ok": False, "error": "playwright not installed"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


# ── Workspace collaboration handler ──

def _handle_workspace_status(**kwargs) -> dict:
    """List workspace files with last-modified agent metadata.

    Useful for multi-agent collaboration: see which files are being
    actively edited and by whom, avoiding conflicts.
    """
    workspace = kwargs.get("path", "workspace")
    if not os.path.isdir(workspace):
        return {"ok": True, "files": [],
                "message": f"Workspace dir '{workspace}' does not exist"}

    files = []
    try:
        for root, dirs, fnames in os.walk(workspace):
            # Skip hidden dirs
            dirs[:] = [d for d in dirs if not d.startswith(".")]
            for fname in fnames:
                if fname.startswith("."):
                    continue
                fpath = os.path.join(root, fname)
                rel = os.path.relpath(fpath, workspace)
                stat = os.stat(fpath)
                entry = {
                    "path": rel,
                    "size": stat.st_size,
                    "modified": stat.st_mtime,
                }
                # Check for agent metadata (written by write_file)
                meta_path = fpath + ".meta"
                if os.path.exists(meta_path):
                    try:
                        with open(meta_path) as f:
                            meta = json.load(f)
                        entry["last_agent"] = meta.get("agent", "unknown")
                        entry["task_id"] = meta.get("task_id", "")
                    except (json.JSONDecodeError, OSError):
                        pass
                files.append(entry)
    except OSError as e:
        return {"ok": False, "error": str(e)}

    # Sort by modification time (most recent first)
    files.sort(key=lambda f: f["modified"], reverse=True)
    return {"ok": True, "files": files[:50], "total": len(files)}


# ── Image understanding handler ──

def _handle_analyze_image(**kwargs) -> dict:
    """Analyze an image using a vision-capable LLM.

    Accepts either a URL or a local file path (auto-encoded to base64).
    Returns the model's description / analysis.
    """
    import base64
    import mimetypes

    image_url = kwargs.get("image_url", "")
    image_path = kwargs.get("image_path", "")
    prompt = kwargs.get("prompt", "Describe this image in detail.")

    if not image_url and not image_path:
        return {"ok": False,
                "error": "Provide either image_url or image_path"}

    # Build the image content block (OpenAI vision API format)
    if image_path:
        abs_path = os.path.abspath(image_path)
        if not os.path.exists(abs_path):
            return {"ok": False, "error": f"File not found: {image_path}"}
        # Read + base64 encode
        mime, _ = mimetypes.guess_type(abs_path)
        if not mime:
            mime = "image/png"
        try:
            with open(abs_path, "rb") as f:
                b64 = base64.b64encode(f.read()).decode()
            image_url = f"data:{mime};base64,{b64}"
        except Exception as e:
            return {"ok": False, "error": f"Failed to read image: {e}"}

    # Build multi-modal message
    messages = [
        {
            "role": "user",
            "content": [
                {
                    "type": "image_url",
                    "image_url": {"url": image_url},
                },
                {
                    "type": "text",
                    "text": prompt,
                },
            ],
        }
    ]

    # Use the resilient LLM adapter to call a vision-capable model
    try:
        import asyncio
        from adapters.llm.resilience import get_llm

        llm = get_llm()
        # Prefer a vision-capable model; fallback to default
        vision_model = os.environ.get(
            "CLEO_VISION_MODEL",
            os.environ.get("CLEO_DEFAULT_MODEL", "gpt-4o"),
        )

        loop = None
        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            pass

        if loop and loop.is_running():
            # We're inside an async context — schedule a task
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor() as pool:
                result = pool.submit(
                    asyncio.run, llm.chat(messages, vision_model)
                ).result(timeout=60)
        else:
            result = asyncio.run(llm.chat(messages, vision_model))

        return {
            "ok": True,
            "model": vision_model,
            "analysis": result,
        }
    except Exception as e:
        return {"ok": False, "error": f"Vision analysis failed: {e}"}


# ── Subagent spawn handler ──

def _handle_spawn_subagent(**kwargs) -> dict:
    """Spawn a child agent to handle a subtask dynamically."""
    description = kwargs.get("description", "")
    parent_id = kwargs.get("parent_id", "")

    if not description:
        return {"ok": False, "error": "description parameter required"}
    if not parent_id:
        return {"ok": False, "error": "parent_id parameter required"}

    try:
        from core.subagent import SubagentRegistry
        from core.task_board import TaskBoard

        board = TaskBoard()
        registry = SubagentRegistry(board)

        child_id = registry.spawn(
            parent_id=parent_id,
            description=description,
            mode=kwargs.get("mode", "run"),
            config_overrides={
                k: v for k, v in {
                    "model": kwargs.get("model"),
                    "skills": kwargs.get("skills"),
                }.items() if v is not None
            },
        )

        return {
            "ok": True,
            "child_task_id": child_id,
            "parent_id": parent_id,
            "message": f"Spawned subagent task {child_id[:8]}",
        }
    except ValueError as e:
        return {"ok": False, "error": str(e)}
    except Exception as e:
        return {"ok": False, "error": f"Spawn failed: {e}"}


# ── A2A delegate handler (Phase 5) ──

def _handle_a2a_delegate(**kwargs) -> dict:
    """Delegate a subtask to an external AI agent via A2A protocol."""
    agent_url = kwargs.get("agent_url", "auto")
    message = kwargs.get("message", "")
    files_str = kwargs.get("files", "")
    timeout = int(kwargs.get("timeout", 120))
    stream = kwargs.get("stream", True)

    if not message:
        return {"ok": False, "error": "message parameter required"}

    try:
        from adapters.a2a.client import A2AClient
        import yaml

        # Load config
        config = {}
        config_path = "config/agents.yaml"
        if os.path.exists(config_path):
            with open(config_path, "r", encoding="utf-8") as f:
                config = yaml.safe_load(f) or {}

        client = A2AClient(config)
        if not client.enabled:
            return {"ok": False,
                    "error": "A2A Client disabled. Set a2a.client.enabled=true"}

        # Parse files
        files = [f.strip() for f in files_str.split(",") if f.strip()] \
            if isinstance(files_str, str) and files_str else \
            (files_str if isinstance(files_str, list) else [])

        # Parse required_skills from message context or kwargs
        required_skills = kwargs.get("required_skills", [])
        if isinstance(required_skills, str):
            required_skills = [s.strip() for s in required_skills.split(",")
                               if s.strip()]

        result = client.send_task(
            agent_url=agent_url,
            message=message,
            files=files,
            required_skills=required_skills,
            timeout=timeout,
            stream=bool(stream),
        )

        _audit_log("a2a_delegate",
                    agent_url=result.agent_url,
                    trust=result.trust_level,
                    status=result.status,
                    duration=result.duration)

        return {
            "ok": result.status == "completed",
            "result": result.to_dict(),
        }

    except Exception as e:
        logger.error("a2a_delegate error: %s", e, exc_info=True)
        return {"ok": False, "error": f"A2A delegation failed: {e}"}


# ══════════════════════════════════════════════════════════════════════════════
#  REGISTRY
# ══════════════════════════════════════════════════════════════════════════════

_BUILTIN_TOOLS: list[Tool] = [
    # ── Web tools ──
    Tool("web_search",
         "Search the web. Supports Brave Search and Perplexity Sonar with auto-fallback. "
         "Results are cached for 15 minutes.",
         {"query": {"type": "string", "description": "Search query", "required": True},
          "count": {"type": "integer", "description": "Number of results (1-20, default 5)", "required": False},
          "freshness": {"type": "string", "description": "Time filter: pd (past day), pw (past week), pm (past month), py (past year)", "required": False},
          "country": {"type": "string", "description": "Country code for results, e.g. US, CN, JP, DE", "required": False},
          "search_lang": {"type": "string", "description": "Search language, e.g. en, zh, ja", "required": False},
          "ui_lang": {"type": "string", "description": "UI language, e.g. en-US, zh-CN", "required": False},
          "provider": {"type": "string", "description": "Force search provider: brave or perplexity (auto-detected by default)", "required": False}},
         _handle_web_search, group="web"),

    Tool("web_fetch",
         "Fetch a URL and extract content. Supports text and markdown extraction modes. "
         "Blocks private/internal hostnames. Results cached for 15 minutes.",
         {"url": {"type": "string", "description": "URL to fetch (must include https://)", "required": True},
          "max_chars": {"type": "integer", "description": "Max chars to return (default 8000)", "required": False},
          "extract_mode": {"type": "string", "description": "Extraction mode: 'text' (plain text, default) or 'markdown' (preserves structure)", "required": False},
          "timeout": {"type": "integer", "description": "Request timeout in seconds (default 15)", "required": False}},
         _handle_web_fetch, group="web"),

    # ── Automation tools ──
    Tool("exec",
         "Execute a shell command (approval-gated). Only allowlisted commands are permitted.",
         {"command": {"type": "string", "description": "Shell command to run", "required": True},
          "timeout": {"type": "integer", "description": "Max seconds (default 120)", "required": False}},
         _handle_exec, group="automation"),

    Tool("cron_list",
         "List all existing scheduled cron jobs.",
         {},
         _handle_cron_list, group="automation"),

    Tool("cron_add",
         "Create a new scheduled job (reminder, periodic task, webhook).",
         {"name": {"type": "string", "description": "Job name/description", "required": True},
          "action": {"type": "string", "description": "Action type: 'task' (create task), 'exec' (run command), 'webhook' (HTTP call)", "required": True},
          "payload": {"type": "string", "description": "Action payload (task description, shell command, or webhook URL)", "required": True},
          "schedule_type": {"type": "string", "description": "Schedule type: 'once' (one-shot), 'interval' (repeat), 'cron' (cron expression)", "required": True},
          "schedule": {"type": "string", "description": "Schedule value: ISO datetime for 'once', seconds for 'interval', cron expression for 'cron' (e.g. '*/5 * * * *')", "required": True}},
         _handle_cron_add, group="automation"),

    Tool("process",
         "List running system processes.",
         {},
         _handle_process_list, group="automation"),

    # ── Skill dependency tools ──
    Tool("check_skill_deps",
         "Check which skill CLI tools are installed and which are missing. "
         "Call with no params to check all, or pass a skill name to check one.",
         {"skill": {"type": "string",
                     "description": "Skill name to check (optional; omit for all)",
                     "required": False}},
         _handle_check_skill_deps, group="skill"),

    Tool("install_skill_cli",
         "Install the CLI tool required by a skill. Auto-detects the best package manager "
         "(brew/go/npm/uv) and runs the install. After install, the binary is automatically "
         "approved for exec. Use check_skill_deps first to see what's missing.",
         {"skill": {"type": "string",
                     "description": "Skill name whose CLI to install (e.g. 'apple-reminders', 'github')",
                     "required": True}},
         _handle_install_skill_cli, group="skill"),

    Tool("search_skills",
         "Search the remote skill registry for new capabilities. Use this when you "
         "need a skill that isn't currently installed. Returns matching skills with "
         "descriptions, versions, and install status.",
         {"query": {"type": "string",
                     "description": "Search query (skill name, description, or tag, e.g. 'pdf', 'browser', 'email')",
                     "required": True},
          "limit": {"type": "integer",
                     "description": "Max results to return (default 10)",
                     "required": False}},
         _handle_search_skills, group="skill"),

    Tool("install_remote_skill",
         "Download and install a skill from the remote registry. After installation, "
         "the skill is immediately available via hot-reload (no restart needed). "
         "CLI dependencies are auto-installed if possible. Use search_skills first.",
         {"slug": {"type": "string",
                    "description": "Skill slug from search_skills results (e.g. 'pdf-rotate', 'browser-control')",
                    "required": True},
          "agent": {"type": "string",
                     "description": "Agent ID to add skill to (default: add to all agents)",
                     "required": False},
          "add_to_all": {"type": "boolean",
                          "description": "Add skill to all agents (default: true)",
                          "required": False}},
         _handle_install_remote_skill, group="skill"),

    # ── Browser automation tools ──
    Tool("browser_navigate",
         "Open a URL in a headless browser. Use for web scraping, form filling, "
         "or interacting with web pages that require JavaScript rendering.",
         {"url": {"type": "string",
                   "description": "URL to navigate to (must include https://)",
                   "required": True}},
         _handle_browser_navigate, group="browser"),

    Tool("browser_click",
         "Click an element on the page by CSS selector.",
         {"selector": {"type": "string",
                        "description": "CSS selector (e.g. 'button.submit', '#login-btn')",
                        "required": True}},
         _handle_browser_click, group="browser"),

    Tool("browser_fill",
         "Fill a form input field with text.",
         {"selector": {"type": "string",
                        "description": "CSS selector for the input field",
                        "required": True},
          "value": {"type": "string",
                     "description": "Text value to enter",
                     "required": True}},
         _handle_browser_fill, group="browser"),

    Tool("browser_get_text",
         "Extract text content from the page or a specific element.",
         {"selector": {"type": "string",
                        "description": "CSS selector (default: 'body' for full page)",
                        "required": False}},
         _handle_browser_get_text, group="browser"),

    Tool("browser_screenshot",
         "Take a screenshot of the current page. Returns file path.",
         {"full_page": {"type": "boolean",
                         "description": "Capture full scrollable page (default: false)",
                         "required": False},
          "selector": {"type": "string",
                        "description": "Screenshot only this element (CSS selector)",
                        "required": False}},
         _handle_browser_screenshot, group="browser"),

    Tool("browser_evaluate",
         "Execute JavaScript code in the page context. Returns the result.",
         {"expression": {"type": "string",
                          "description": "JavaScript expression to evaluate",
                          "required": True}},
         _handle_browser_evaluate, group="browser"),

    Tool("browser_page_info",
         "Get current browser page info (URL, title).",
         {},
         _handle_browser_page_info, group="browser"),

    # ── Media tools ──
    Tool("screenshot",
         "Capture a screenshot of the current desktop (macOS only).",
         {},
         _handle_screenshot, group="media"),

    Tool("notify",
         "Send a desktop notification (macOS only).",
         {"title": {"type": "string", "description": "Notification title", "required": True},
          "message": {"type": "string", "description": "Notification body", "required": False}},
         _handle_notify, group="media"),

    Tool("transcribe",
         "Transcribe an audio file to text using OpenAI Whisper API. "
         "Supports mp3, mp4, m4a, wav, webm, ogg, flac (max 25 MB). "
         "Requires OPENAI_API_KEY.",
         {"file_path": {"type": "string", "description": "Path to the audio file", "required": True},
          "language": {"type": "string",
                       "description": "ISO 639-1 language code hint (e.g. 'zh', 'en', 'ja'). "
                                      "Optional — Whisper auto-detects if omitted.",
                       "required": False}},
         _handle_transcribe, group="media",
         requires_env=["OPENAI_API_KEY"]),

    Tool("tts",
         "Text-to-speech: synthesize text into an audio file. "
         "Multi-provider with automatic fallback (OpenAI/ElevenLabs/MiniMax/local). "
         "Returns path to the generated audio file.",
         {"text": {"type": "string", "description": "Text to synthesize (max ~5000 chars)", "required": True},
          "voice": {"type": "string",
                    "description": "Voice ID — provider-specific. "
                    "OpenAI: alloy/echo/nova/shimmer/fable/onyx. "
                    "ElevenLabs: rachel/adam/sam/josh/bella. "
                    "MiniMax: presenter_male/presenter_female/female-shaonv.",
                    "required": False},
          "speed": {"type": "number", "description": "Speed multiplier 0.5-2.0 (default 1.0)", "required": False},
          "provider": {"type": "string", "description": "Force provider: openai/elevenlabs/minimax/local", "required": False},
          "output_format": {"type": "string", "description": "Output format: mp3/wav/ogg (default mp3)", "required": False}},
         _handle_tts, group="media"),

    Tool("list_voices",
         "List available TTS voices across all configured providers.",
         {"provider": {"type": "string", "description": "Filter by provider name (optional)", "required": False}},
         _handle_list_voices, group="media"),

    # ── Filesystem tools ──
    Tool("read_file",
         "Read a file from the project directory.",
         {"path": {"type": "string", "description": "File path relative to project root", "required": True},
          "max_lines": {"type": "integer", "description": "Max lines to read (default 200)", "required": False}},
         _handle_read_file, group="fs"),

    Tool("write_file",
         "Write content to a file in the project directory.",
         {"path": {"type": "string", "description": "File path relative to project root", "required": True},
          "content": {"type": "string", "description": "Content to write", "required": True}},
         _handle_write_file, group="fs"),

    Tool("list_dir",
         "List directory contents.",
         {"path": {"type": "string", "description": "Directory path (default: project root)", "required": False}},
         _handle_list_dir, group="fs"),

    Tool("edit_file",
         "Find-and-replace edit in a project file. The old_str must be unique in the file.",
         {"path": {"type": "string", "description": "File path relative to project root", "required": True},
          "old_str": {"type": "string", "description": "Exact text to find (must be unique)", "required": True},
          "new_str": {"type": "string", "description": "Replacement text", "required": True}},
         _handle_edit_file, group="fs"),

    # ── Memory tools ──
    Tool("memory_search",
         "Search episodic memory for past problem→solution cases.",
         {"query": {"type": "string", "description": "Search query", "required": True},
          "limit": {"type": "integer", "description": "Max results (default 5)", "required": False}},
         _handle_memory_search, group="memory"),

    Tool("memory_save",
         "Save a problem→solution case to episodic memory for future recall.",
         {"problem": {"type": "string", "description": "Problem description", "required": True},
          "solution": {"type": "string", "description": "Solution description", "required": True},
          "tags": {"type": "string", "description": "Comma-separated tags", "required": False}},
         _handle_memory_save, group="memory"),

    Tool("kb_search",
         "Search the shared knowledge base for notes and insights.",
         {"query": {"type": "string", "description": "Search query", "required": True},
          "limit": {"type": "integer", "description": "Max results (default 5)", "required": False}},
         _handle_kb_search, group="memory"),

    Tool("kb_write",
         "Create or update a note in the shared knowledge base (Zettelkasten).",
         {"topic": {"type": "string", "description": "Note topic/title", "required": True},
          "content": {"type": "string", "description": "Note content", "required": True},
          "tags": {"type": "string", "description": "Comma-separated tags", "required": False}},
         _handle_kb_write, group="memory"),

    # ── Task tools ──
    Tool("task_create",
         "Create a new task on the task board.",
         {"description": {"type": "string", "description": "Task description", "required": True}},
         _handle_task_create, group="task"),

    Tool("task_status",
         "Get task status. Without task_id, lists recent tasks.",
         {"task_id": {"type": "string", "description": "Task ID or prefix (optional)", "required": False}},
         _handle_task_status, group="task"),

    # ── Messaging tools ──
    Tool("send_mail",
         "Send a message to another agent's mailbox for inter-agent communication.",
         {"to": {"type": "string", "description": "Target agent ID", "required": True},
          "content": {"type": "string", "description": "Message content", "required": True},
          "msg_type": {"type": "string", "description": "Message type (default: message)", "required": False}},
         _handle_send_mail, group="messaging"),

    Tool("generate_doc",
         "Generate a document file from content. "
         "Supports 8 formats: pdf, docx/word, xlsx/excel, pptx/powerpoint, csv, txt, md, html. "
         "CJK (Chinese/Japanese/Korean) fully supported. "
         "PDF auto-falls-back to DOCX on failure. "
         "After generating, use send_file to deliver to the user.",
         {"format": {"type": "string",
                     "description": "Output format: 'pdf', 'docx'/'word', 'xlsx'/'excel', "
                                    "'pptx'/'powerpoint', 'csv', 'txt', 'md', 'html'",
                     "required": True},
          "content": {"type": "string",
                      "description": "Document content (IMPORTANT: keep under 2000 chars to avoid truncation; "
                                     "use concise formatting). Markdown text for pdf/docx/pptx/html/txt, "
                                     "JSON array or markdown table for xlsx/csv",
                      "required": True},
          "title": {"type": "string",
                    "description": "Document title (optional)",
                    "required": False},
          "output_path": {"type": "string",
                          "description": "Output file path (default: /tmp/doc_<ts>.<ext>)",
                          "required": False}},
         _handle_generate_doc, group="fs"),

    Tool("send_file",
         "Send a file to the user via their chat channel (Telegram/Discord/Feishu/Slack). "
         "Use this after creating a file (with generate_doc or write_file) to deliver it to the user. "
         "Only works when the task originates from a channel message.",
         {"file_path": {"type": "string",
                        "description": "Absolute or relative path to the file to send",
                        "required": True},
          "caption": {"type": "string",
                      "description": "Optional caption/message to include with the file",
                      "required": False}},
         _handle_send_file, group="messaging"),

    Tool("message",
         "Send a message to the user via their chat channel. "
         "Can send text, a file, or both. Use for proactive communication "
         "(progress updates, questions, delivering results). "
         "Prefer this over send_file when you also want to include a text message.",
         {"text": {"type": "string",
                   "description": "Text message to send to the user",
                   "required": False},
          "file_path": {"type": "string",
                        "description": "Path to file to send (optional)",
                        "required": False},
          "caption": {"type": "string",
                      "description": "Caption for the file (only used when file_path is set)",
                      "required": False}},
         _handle_message, group="messaging"),

    # ── Collaboration tools ──
    Tool("workspace_status",
         "List workspace files with modification info and which agent last touched "
         "each file. Useful for avoiding edit conflicts in multi-agent collaboration.",
         {"path": {"type": "string",
                   "description": "Workspace directory path (default: 'workspace')",
                   "required": False}},
         _handle_workspace_status, group="fs"),

    # ── Vision tools ──
    Tool("analyze_image",
         "Analyze an image using a vision-capable LLM. Accepts a URL or local file path. "
         "Use this to understand screenshots, diagrams, photos, or any visual content.",
         {"image_url": {"type": "string",
                        "description": "URL of the image to analyze (HTTPS or data: URI)",
                        "required": False},
          "image_path": {"type": "string",
                         "description": "Local file path to the image (auto-encoded to base64)",
                         "required": False},
          "prompt": {"type": "string",
                     "description": "What to analyze or ask about the image (default: 'Describe this image in detail.')",
                     "required": False}},
         _handle_analyze_image, group="media"),

    # ── Subagent tools ──
    Tool("spawn_subagent",
         "Dynamically spawn a child agent to handle a subtask. The child runs as an "
         "independent task and notifies you upon completion via mailbox.",
         {"description": {"type": "string",
                          "description": "Task description for the child agent",
                          "required": True},
          "parent_id": {"type": "string",
                        "description": "Your agent ID (the parent spawning this child)",
                        "required": True},
          "mode": {"type": "string",
                   "description": "Spawn mode: 'run' (one-shot) or 'session' (persistent)",
                   "required": False},
          "model": {"type": "string",
                    "description": "Override LLM model for the child agent",
                    "required": False},
          "skills": {"type": "string",
                     "description": "Comma-separated skill names for the child",
                     "required": False}},
         _handle_spawn_subagent, group="task"),

    # ── A2A delegation tool (Phase 5) ──
    Tool("a2a_delegate",
         "Delegate a subtask to an external AI agent via the A2A (Agent-to-Agent) protocol. "
         "Use when the task requires capabilities Cleo doesn't have (chart generation, "
         "specialized data analysis, image generation, etc). Set agent_url='auto' to "
         "automatically find the best matching agent by required_skills.",
         {"agent_url": {"type": "string",
                        "description": "Target agent URL or 'auto' for automatic matching",
                        "required": True},
          "message": {"type": "string",
                      "description": "Task description for the external agent (English preferred)",
                      "required": True},
          "files": {"type": "string",
                    "description": "Comma-separated file paths to attach (verified agents only)",
                    "required": False},
          "required_skills": {"type": "string",
                              "description": "Comma-separated skill tags for auto-matching",
                              "required": False},
          "timeout": {"type": "integer",
                      "description": "Max wait seconds (default 120)",
                      "required": False}},
         _handle_a2a_delegate, group="a2a"),
]

# Keyed registry for fast lookup
_registry: dict[str, Tool] = {t.name: t for t in _BUILTIN_TOOLS}


def get_tool(name: str) -> Tool | None:
    """Get a tool by name."""
    return _registry.get(name)


def list_all_tools() -> list[Tool]:
    """List all registered tools."""
    return list(_BUILTIN_TOOLS)


def get_available_tools(agent_config: dict | None = None) -> list[Tool]:
    """Get tools available to an agent based on its config.

    agent_config can have:
      tools.profile: "minimal" | "coding" | "full"
      tools.allow: ["tool_name", "group:web"]
      tools.deny: ["tool_name", "group:automation"]
    """
    tools_cfg = (agent_config or {}).get("tools", {})
    profile = tools_cfg.get("profile", "full")
    allow_extra = set(tools_cfg.get("allow", []))
    deny_list = set(tools_cfg.get("deny", []))

    # Expand groups
    expanded_allow = set()
    for item in allow_extra:
        if item in TOOL_GROUPS:
            expanded_allow.update(TOOL_GROUPS[item])
        else:
            expanded_allow.add(item)

    expanded_deny = set()
    for item in deny_list:
        if item in TOOL_GROUPS:
            expanded_deny.update(TOOL_GROUPS[item])
        else:
            expanded_deny.add(item)

    # Base set from profile
    base_set = TOOL_PROFILES.get(profile)  # None = full

    available = []
    for tool in _BUILTIN_TOOLS:
        # Check deny first (deny always wins)
        if tool.name in expanded_deny:
            continue

        # Check base profile
        if base_set is not None and tool.name not in base_set:
            # Not in base profile — check if explicitly allowed
            if tool.name not in expanded_allow:
                continue

        # Check env requirements
        if not tool.is_available():
            continue

        available.append(tool)

    return available


def build_tools_prompt(agent_config: dict | None = None) -> str:
    """Build the tools section for agent system prompt."""
    tools = get_available_tools(agent_config)
    if not tools:
        return ""

    lines = [
        "## Available Tools",
        "",
        "You can invoke tools by including a JSON block in your response.",
        "Use this EXACT format (any of these formats work):",
        "",
        "Format 1 (preferred):",
        "```tool",
        '{"tool": "tool_name", "params": {"param1": "value1"}}',
        "```",
        "",
        "Format 2:",
        "<tool_code>",
        '{"tool": "tool_name", "params": {"param1": "value1"}}',
        "</tool_code>",
        "",
        "IMPORTANT: Use ONE tool per block. The tool JSON must have a \"tool\" key and a \"params\" key.",
        "After tool execution, you will receive the results and can continue.",
        "",
        "Available tools:",
        "",
    ]
    for t in tools:
        lines.append(t.to_prompt())

    return "\n".join(lines)


def build_tools_schemas(agent_config: dict | None = None) -> list[dict]:
    """Build tool schemas for function-calling LLMs."""
    tools = get_available_tools(agent_config)
    return [t.to_schema() for t in tools]


# ══════════════════════════════════════════════════════════════════════════════
#  V0.02 TOOL SCOPE — Category-based tool loading
# ══════════════════════════════════════════════════════════════════════════════

# Base tools always loaded (memory + messaging essentials)
_BASE_TOOL_NAMES: set[str] = {
    "memory_search", "memory_save", "kb_search", "kb_write",
    "send_mail", "send_file", "message",
}

# Startup validation: ensure all base tool names exist in _BUILTIN_TOOLS
_builtin_names = {t.name for t in _BUILTIN_TOOLS}
_missing_base = _BASE_TOOL_NAMES - _builtin_names
if _missing_base:
    logger.error("ToolScope: _BASE_TOOL_NAMES references non-existent tools: %s", _missing_base)

# Mapping from SubTaskSpec tool_hint values → TOOL_GROUPS keys
_HINT_TO_GROUP: dict[str, str] = {
    "web":       "group:web",
    "fs":        "group:fs",
    "automation": "group:automation",
    "media":     "group:media",
    "browser":   "group:browser",
    "memory":    "group:memory",
    "messaging": "group:messaging",
    "task":      "group:task",
    "skill":     "group:skill",
    "a2a_delegate": "group:a2a",
}


def get_scoped_tools(tool_hints: list[str],
                     agent_config: dict | None = None) -> list[Tool]:
    """Get tools scoped to specific categories (V0.02 ToolScope).

    Loads base tools (memory + messaging) plus tools matching tool_hints.
    Falls back to full profile if tool_hints is empty.

    Args:
        tool_hints: List of category hints (e.g. ["web", "fs"]).
        agent_config: Agent config dict with tools.deny etc.

    Returns:
        Scoped tool list (typically 9-14 tools vs 33 in full coding profile).
    """
    if not tool_hints:
        # Fallback: V0.01 behavior — load full profile
        return get_available_tools(agent_config)

    # Build allowed tool names: base + categories from hints
    allowed_names: set[str] = set(_BASE_TOOL_NAMES)
    for hint in tool_hints:
        group_key = _HINT_TO_GROUP.get(hint)
        if group_key and group_key in TOOL_GROUPS:
            allowed_names.update(TOOL_GROUPS[group_key])

    # Apply deny list from agent config
    tools_cfg = (agent_config or {}).get("tools", {})
    deny_list = set(tools_cfg.get("deny", []))
    expanded_deny: set[str] = set()
    for item in deny_list:
        if item in TOOL_GROUPS:
            expanded_deny.update(TOOL_GROUPS[item])
        else:
            expanded_deny.add(item)

    available = []
    for tool in _BUILTIN_TOOLS:
        if tool.name not in allowed_names:
            continue
        if tool.name in expanded_deny:
            continue
        if not tool.is_available():
            continue
        available.append(tool)

    return available


def build_scoped_tools_prompt(tool_hints: list[str],
                              agent_config: dict | None = None) -> str:
    """Build tools prompt section for scoped tool set."""
    tools = get_scoped_tools(tool_hints, agent_config)
    if not tools:
        return ""

    lines = [
        "## Available Tools",
        "",
        "You can invoke tools by including a JSON block in your response.",
        "Use this EXACT format (any of these formats work):",
        "",
        "Format 1 (preferred):",
        "```tool",
        '{"tool": "tool_name", "params": {"param1": "value1"}}',
        "```",
        "",
        "Format 2:",
        "<tool_code>",
        '{"tool": "tool_name", "params": {"param1": "value1"}}',
        "</tool_code>",
        "",
        "IMPORTANT: Use ONE tool per block. The tool JSON must have a \"tool\" key and a \"params\" key.",
        "After tool execution, you will receive the results and can continue.",
        "",
        f"Available tools (scoped to: {', '.join(tool_hints)}):",
        "",
    ]
    for t in tools:
        lines.append(t.to_prompt())

    return "\n".join(lines)


def build_scoped_tools_schemas(tool_hints: list[str],
                               agent_config: dict | None = None) -> list[dict]:
    """Build tool schemas for scoped tool set (function-calling LLMs)."""
    tools = get_scoped_tools(tool_hints, agent_config)
    return [t.to_schema() for t in tools]


# ══════════════════════════════════════════════════════════════════════════════
#  TOOL INVOCATION PARSER
# ══════════════════════════════════════════════════════════════════════════════

# Matches ```tool\n{...}\n``` blocks in agent output
_TOOL_BLOCK_RE = re.compile(
    r'```tool\s*\n(\{[^`]+?\})\s*\n```',
    re.DOTALL)

# Fallback: matches <tool_code>\n{...}\n</tool_code> blocks (Minimax format)
_TOOL_CODE_RE = re.compile(
    r'<tool_code>\s*\n?([\s\S]+?)\n?\s*</tool_code>',
    re.DOTALL)

# Fallback: matches ```json\n{"tool":...}\n``` blocks
_JSON_BLOCK_RE = re.compile(
    r'```(?:json)?\s*\n(\{"tool"\s*:\s*[^`]+?\})\s*\n```',
    re.DOTALL)

# Fallback: matches <minimax:tool_call>...</ or <invoke name="tool", ...>
_MINIMAX_CALL_RE = re.compile(
    r'<(?:minimax:tool_call|invoke)\b[^>]*>([\s\S]+?)</(?:minimax:tool_call|tool_code|invoke)>',
    re.DOTALL)

# Fallback: extract from <invoke name="tool_name", "params": {...}>
_INVOKE_ATTR_RE = re.compile(
    r'<invoke\s+name\s*=\s*"(\w+)"[^>]*?(?:"params"\s*:\s*(\{[^}]*\}))?',
    re.DOTALL)


def _try_parse_arrow_syntax(raw: str) -> dict | None:
    """Parse Minimax-style arrow syntax to JSON.

    Handles formats like:
        { tool => 'web_search', args => { --query "hello" } }
        { tool => 'web_search', params => { query: "hello" } }
    """
    try:
        # Strip outer braces and normalize
        s = raw.strip()
        if not s.startswith("{"):
            return None

        # Extract tool name: tool => 'name' or tool => "name"
        tool_match = re.search(r'''tool\s*(?:=>|:)\s*['"](\w+)['"]''', s)
        if not tool_match:
            return None
        tool_name = tool_match.group(1)

        # Extract params/args block
        params = {}
        args_match = re.search(
            r'(?:args|params)\s*(?:=>|:)\s*\{([^}]*)\}', s, re.DOTALL)
        if args_match:
            args_raw = args_match.group(1)
            # Parse --key "value" patterns (CLI-style)
            for m in re.finditer(
                    r'--(\w+)\s+["\']([^"\']*)["\']', args_raw):
                params[m.group(1)] = m.group(2)
            # Parse key: "value" or key => "value" patterns
            for m in re.finditer(
                    r'(\w+)\s*(?:=>|:)\s*["\']([^"\']*)["\']', args_raw):
                params[m.group(1)] = m.group(2)
            # Parse key: number patterns
            for m in re.finditer(
                    r'(\w+)\s*(?:=>|:)\s*(\d+)', args_raw):
                params[m.group(1)] = int(m.group(2))

        return {"tool": tool_name, "params": params}
    except Exception:
        return None


def parse_tool_calls(text: str) -> list[dict]:
    """Extract tool invocation blocks from agent output.

    Supports multiple formats for LLM compatibility:
    1. Standard: ```tool\n{"tool":"name","params":{...}}\n```
    2. Minimax: <tool_code>\n{tool=>'name',args=>{...}}\n</tool_code>
    3. JSON block: ```json\n{"tool":"name","params":{...}}\n```

    Returns list of {"tool": "name", "params": {...}}
    """
    calls = []

    # 1. Standard format (```tool ... ```)
    for match in _TOOL_BLOCK_RE.finditer(text):
        try:
            data = json.loads(match.group(1))
            if "tool" in data:
                calls.append({
                    "tool": data["tool"],
                    "params": data.get("params", {}),
                    "raw": match.group(0),
                })
        except json.JSONDecodeError:
            continue

    if calls:
        return calls

    # 2. <tool_code> blocks (Minimax arrow syntax)
    for match in _TOOL_CODE_RE.finditer(text):
        raw_content = match.group(1).strip()
        # Try JSON first
        try:
            data = json.loads(raw_content)
            if "tool" in data:
                calls.append({
                    "tool": data["tool"],
                    "params": data.get("params", data.get("args", {})),
                    "raw": match.group(0),
                })
                continue
        except json.JSONDecodeError:
            pass
        # Try arrow syntax
        parsed = _try_parse_arrow_syntax(raw_content)
        if parsed:
            calls.append({
                "tool": parsed["tool"],
                "params": parsed.get("params", {}),
                "raw": match.group(0),
            })

    if calls:
        return calls

    # 3. ```json blocks with tool key
    for match in _JSON_BLOCK_RE.finditer(text):
        try:
            data = json.loads(match.group(1))
            if "tool" in data:
                calls.append({
                    "tool": data["tool"],
                    "params": data.get("params", data.get("args", {})),
                    "raw": match.group(0),
                })
        except json.JSONDecodeError:
            continue

    if calls:
        return calls

    # 4. <minimax:tool_call> or <invoke name="..."> blocks
    for match in _MINIMAX_CALL_RE.finditer(text):
        raw_content = match.group(1).strip()
        # Try JSON parse
        try:
            data = json.loads(raw_content)
            if "tool" in data:
                calls.append({
                    "tool": data["tool"],
                    "params": data.get("params", data.get("args", {})),
                    "raw": match.group(0),
                })
                continue
        except json.JSONDecodeError:
            pass
        # Try arrow syntax
        parsed = _try_parse_arrow_syntax(raw_content)
        if parsed:
            calls.append({
                "tool": parsed["tool"],
                "params": parsed.get("params", {}),
                "raw": match.group(0),
            })

    if calls:
        return calls

    # 4b. <invoke name="tool_name", "params": {...}> (attribute-style)
    for match in _INVOKE_ATTR_RE.finditer(text):
        tool_name = match.group(1)
        params = {}
        if match.group(2):
            try:
                params = json.loads(match.group(2))
            except json.JSONDecodeError:
                pass
        calls.append({
            "tool": tool_name,
            "params": params,
            "raw": match.group(0),
        })

    if not calls and any(kw in text for kw in
                         ["web_search", "web_fetch", "exec", "read_file",
                          "write_file", "memory_search"]):
        logger.warning("Tool keywords found in LLM output but no parseable "
                       "tool blocks detected. First 300 chars: %s",
                       text[:300].replace("\n", "\\n"))

    return calls


# ══════════════════════════════════════════════════════════════════════════════
#  PARAMETER SANITIZATION — defence-in-depth for LLM-generated params
# ══════════════════════════════════════════════════════════════════════════════

# Sensitive files that tools should never read or write
_SENSITIVE_FILENAMES = {
    ".env", ".env.local", ".env.production", ".env.development",
    "agents.yaml", "exec_approvals.json", "chain_contracts.json",
    ".git/config", ".netrc", ".npmrc", ".pypirc",
    "id_rsa", "id_ed25519", "authorized_keys",
}

# Sensitive path fragments (blocked anywhere in path)
_SENSITIVE_PATH_FRAGMENTS = {
    ".ssh", ".gnupg", ".aws", ".config/gcloud",
}

# Tools whose "path" param must be checked
_FS_TOOLS = {"read_file", "write_file", "edit_file", "list_dir"}

# Tools whose "url" param must be scheme-checked
_NET_TOOLS = {"web_fetch", "web_search"}


def sanitize_params(tool_name: str, params: dict,
                    tool: "Tool | None" = None) -> dict | str:
    """Validate and sanitize LLM-generated tool parameters.

    Returns sanitized params dict on success, or error string on rejection.
    Checks performed:
      1. Type coercion — cast to schema-declared types
      2. Path safety  — block sensitive files, enforce project scope
      3. URL safety   — enforce https, block private IPs (defence-in-depth)
    """
    if not isinstance(params, dict):
        return "Parameters must be a JSON object"

    # ── 1. Type coercion ──
    if tool and tool.parameters:
        for pname, pinfo in tool.parameters.items():
            if pname not in params:
                continue
            expected = pinfo.get("type", "string")
            val = params[pname]
            try:
                if expected == "integer" and not isinstance(val, int):
                    params[pname] = int(val)
                elif expected == "number" and not isinstance(val, (int, float)):
                    params[pname] = float(val)
                elif expected == "boolean" and not isinstance(val, bool):
                    params[pname] = str(val).lower() in ("true", "1", "yes")
                elif expected == "string" and not isinstance(val, str):
                    params[pname] = str(val)
            except (ValueError, TypeError):
                return f"Parameter '{pname}' must be {expected}, got {type(val).__name__}"

    # ── 2. Path safety (filesystem tools) ──
    if tool_name in _FS_TOOLS:
        raw_path = params.get("path", "")
        if not isinstance(raw_path, str) or not raw_path.strip():
            return "Missing or empty 'path' parameter"

        # Normalise to block encoded traversal  (e.g. %2e%2e)
        decoded_path = urllib.parse.unquote(raw_path)

        # Block null bytes (can bypass os.path checks)
        if "\x00" in decoded_path:
            return "Null bytes not allowed in path"

        # Block sensitive filenames
        basename = os.path.basename(decoded_path)
        if basename.lower() in _SENSITIVE_FILENAMES:
            return f"Access to sensitive file '{basename}' is blocked"

        # Block sensitive path fragments
        norm = os.path.normpath(decoded_path).replace("\\", "/").lower()
        for frag in _SENSITIVE_PATH_FRAGMENTS:
            if frag in norm:
                return f"Path contains blocked segment '{frag}'"

        # Write-specific: block hidden dotfiles at project root
        if tool_name == "write_file" and basename.startswith("."):
            return f"Cannot write to hidden file '{basename}' (dotfiles are protected)"

        # Replace raw path with decoded version for consistency
        params["path"] = decoded_path

    # ── 3. URL safety (network tools) ──
    if tool_name in _NET_TOOLS and "url" in params:
        url = params.get("url", "")
        if not isinstance(url, str):
            return "URL must be a string"

        # Enforce https (allow http only for localhost dev, but that's
        # already blocked by _is_private_hostname)
        parsed = urllib.parse.urlparse(url)
        if parsed.scheme not in ("https", "http"):
            return f"URL scheme '{parsed.scheme}' not allowed — use https://"

        # Defence-in-depth: re-check private hostnames at param level
        # (web_fetch already checks, but belt-and-suspenders)
        hostname = parsed.hostname or ""
        if _is_private_hostname(hostname):
            return f"Blocked: private/internal hostname '{hostname}'"

    return params


def execute_tool_calls(calls: list[dict],
                       agent_config: dict | None = None) -> list[dict]:
    """Execute parsed tool calls and return results.

    Returns list of {"tool": "name", "result": {...}}
    """
    available = {t.name for t in get_available_tools(agent_config)}
    results = []
    for call in calls:
        name = call["tool"]
        if name not in available:
            results.append({
                "tool": name,
                "result": {"ok": False, "error": f"Tool '{name}' not available"},
            })
            continue

        tool = _registry.get(name)
        if not tool:
            results.append({
                "tool": name,
                "result": {"ok": False, "error": f"Unknown tool: {name}"},
            })
            continue

        # ── Sanitize parameters before execution ──
        raw_params = call.get("params", {})
        sanitized = sanitize_params(name, dict(raw_params), tool)
        if isinstance(sanitized, str):
            # sanitize_params returned an error message
            logger.warning("Tool %s params rejected: %s (raw: %s)",
                           name, sanitized, str(raw_params)[:200])
            results.append({
                "tool": name,
                "result": {"ok": False, "error": f"Parameter validation: {sanitized}"},
            })
            continue

        logger.info("Executing tool: %s(%s)", name,
                     str(sanitized)[:100])
        result = tool.execute(**sanitized)
        results.append({"tool": name, "result": result})

    return results
